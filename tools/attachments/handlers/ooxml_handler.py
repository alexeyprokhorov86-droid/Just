"""
OOXML handler — docx/xlsx/pptx по содержимому ZIP-архива.

Конкретный формат определяется по filename либо по структуре ZIP.
"""
from __future__ import annotations

import io
import logging

from .._prompts import build_analysis_prompt

log = logging.getLogger("tools.attachments.ooxml")


def _detect_ooxml_kind(file_bytes: bytes, filename: str) -> str:
    """'docx' | 'xlsx' | 'pptx' | 'unknown'."""
    low = filename.lower()
    if low.endswith(".docx"):
        return "docx"
    if low.endswith((".xlsx", ".xls", ".xlsm")):
        return "xlsx"
    if low.endswith((".pptx", ".ppt")):
        return "pptx"
    # Фоллбек через чтение структуры ZIP
    try:
        import zipfile
        with zipfile.ZipFile(io.BytesIO(file_bytes)) as z:
            names = z.namelist()
            if any(n.startswith("word/") for n in names):
                return "docx"
            if any(n.startswith("xl/") for n in names):
                return "xlsx"
            if any(n.startswith("ppt/") for n in names):
                return "pptx"
    except Exception:
        pass
    return "unknown"


def _extract_docx(file_bytes: bytes) -> str:
    try:
        from docx import Document
    except ImportError:
        return ""
    try:
        doc = Document(io.BytesIO(file_bytes))
        parts: list[str] = []
        for para in doc.paragraphs:
            if para.text.strip():
                parts.append(para.text)
        for table in doc.tables:
            rows = []
            for row in table.rows:
                cells = [cell.text.strip() for cell in row.cells]
                if any(cells):
                    rows.append(" | ".join(cells))
            if rows:
                parts.append("\n" + "\n".join(rows))
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"docx extract failed: {e}")
        return ""


def _extract_xlsx(file_bytes: bytes) -> str:
    try:
        import openpyxl
    except ImportError:
        return ""
    try:
        wb = openpyxl.load_workbook(io.BytesIO(file_bytes), data_only=True, read_only=True)
        parts: list[str] = []
        for sheet_name in wb.sheetnames:
            ws = wb[sheet_name]
            parts.append(f"=== Лист: {sheet_name} ===")
            rows_out: list[str] = []
            for row in ws.iter_rows(values_only=True, max_row=200):
                row_text = " | ".join("" if v is None else str(v) for v in row)
                if row_text.strip(" |"):
                    rows_out.append(row_text)
            parts.extend(rows_out)
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"xlsx extract failed: {e}")
        return ""


def _extract_pptx(file_bytes: bytes) -> str:
    try:
        from pptx import Presentation
    except ImportError:
        return ""
    try:
        prs = Presentation(io.BytesIO(file_bytes))
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, 1):
            parts.append(f"=== Слайд {i} ===")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n".join(parts)
    except Exception as e:
        log.warning(f"pptx extract failed: {e}")
        return ""


def analyze_ooxml(
    *,
    file_bytes: bytes,
    filename: str,
    chat_context: str,
    gpt_client,
    company_profile: str,
    model: str = "openai/gpt-4.1",
) -> dict:
    kind = _detect_ooxml_kind(file_bytes, filename)
    errors: list[str] = []

    if kind == "docx":
        extracted = _extract_docx(file_bytes)
        label = "Word документ"
    elif kind == "xlsx":
        extracted = _extract_xlsx(file_bytes)
        label = "Excel документ"
    elif kind == "pptx":
        extracted = _extract_pptx(file_bytes)
        label = "PowerPoint документ"
    else:
        return {
            "document_type": "unknown_zip",
            "extracted_text": "",
            "structured_fields": {"kind": "unknown"},
            "summary": "",
            "confidence": 0.0,
            "errors": ["Не удалось определить тип OOXML"],
        }

    summary = ""
    if gpt_client is not None and extracted:
        prompt = build_analysis_prompt(
            company_profile=company_profile,
            doc_type=label,
            doc_content=extracted[:20000],
            chat_context=chat_context,
            filename=filename,
        )
        try:
            response = gpt_client.chat.completions.create(
                model=model,
                max_tokens=2000,
                messages=[{"role": "user", "content": prompt}],
            )
            summary = (response.choices[0].message.content or "").strip()
        except Exception as e:
            errors.append(f"LLM analysis failed: {e}")

    return {
        "document_type": kind,
        "extracted_text": extracted,
        "structured_fields": {"ooxml_kind": kind},
        "summary": summary,
        "confidence": 1.0 if summary and not errors else (0.3 if extracted else 0.0),
        "errors": errors,
    }
