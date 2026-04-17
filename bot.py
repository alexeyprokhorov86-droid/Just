"""
Telegram Bot для логирования сообщений из групповых чатов.
Версия 2.0 - с ролями пользователей и расширенным контекстом.

Функции:
- Логирование всех сообщений в PostgreSQL
- Анализ изображений, PDF, Excel, Word, PowerPoint через Claude Vision
- Учёт ролей пользователей
- Контекст чата за 8 дней с учётом связанных сообщений
"""

import os
import re
import json
import logging
import base64
import threading
import datetime
from datetime import datetime, timedelta
from telegram import Update, InlineKeyboardButton, InlineKeyboardMarkup
from telegram.ext import Application, MessageHandler, CommandHandler, filters, ContextTypes
import psycopg2
from psycopg2 import sql
from dotenv import load_dotenv
from openai import OpenAI
from telegram.ext import CallbackQueryHandler
from rag_agent import process_rag_query, index_new_message
from telegram.helpers import escape_markdown
from apscheduler.schedulers.asyncio import AsyncIOScheduler
from apscheduler.triggers.cron import CronTrigger
from company_context import get_company_profile
from fact_extractor import extract_and_save_facts
from notifications import get_notify_conversation_handler, handle_ack, notify_status, notify_remind

# Загружаем переменные окружения
# Ищем .env в директории скрипта или в текущей директории
import pathlib
env_path = pathlib.Path(__file__).parent / '.env'
load_dotenv(dotenv_path=env_path if env_path.exists() else None)

# Настройка логирования
logging.basicConfig(
    format='%(asctime)s - %(name)s - %(levelname)s - %(message)s',
    level=logging.INFO
)
logger = logging.getLogger(__name__)

# ============================================================
# КОНФИГУРАЦИЯ
# ============================================================

BOT_TOKEN = os.getenv("BOT_TOKEN")


# Подключение к БД
DB_HOST = os.getenv("DB_HOST", "localhost")
DB_PORT = os.getenv("DB_PORT", "5432")
DB_NAME = os.getenv("DB_NAME", "knowledge_base")
DB_USER = os.getenv("DB_USER", "knowledge")
DB_PASSWORD = os.getenv("DB_PASSWORD")

# ID администратора для запросов ролей (твой Telegram ID)
ADMIN_USER_ID = int(os.getenv("ADMIN_USER_ID", "0"))

# Название группы для отложенного анализа документов
DELAYED_ANALYSIS_CHAT = "Торты Отгрузки"

# Инициализация GPT клиента через RouterAI
ROUTERAI_API_KEY = os.getenv("ROUTERAI_API_KEY")
ROUTERAI_BASE_URL = os.getenv("ROUTERAI_BASE_URL", "https://routerai.ru/api/v1")

gpt_client = None
if ROUTERAI_API_KEY:
    gpt_client = OpenAI(api_key=ROUTERAI_API_KEY, base_url=ROUTERAI_BASE_URL, timeout=400)
    logger.info("GPT-4.1 через RouterAI активирован")
else:
    logger.warning("ROUTERAI_API_KEY не установлен - анализ документов отключён")

# Хранение состояния для назначения ролей
pending_role_assignments = {}

# S3 конфигурация
S3_BUCKET = os.getenv("ATTACHMENTS_BUCKET_NAME", "")
S3_ENDPOINT = os.getenv("ATTACHMENTS_BUCKET_ENDPOINT", "")
S3_REGION = os.getenv("ATTACHMENTS_BUCKET_REGION", "ru-central-1")
S3_ACCESS_KEY = os.getenv("ATTACHMENTS_BUCKET_ACCESS_KEY", "")
S3_SECRET_KEY = os.getenv("ATTACHMENTS_BUCKET_SECRET_KEY", "")
S3_FORCE_PATH = os.getenv("ATTACHMENTS_BUCKET_FORCE_PATH_STYLE", "true").lower() == "true"
 
_s3_client = None
 
 
def get_s3_client():
    """Ленивая инициализация S3 клиента."""
    global _s3_client
    if _s3_client is None and S3_BUCKET and S3_ACCESS_KEY:
        try:
            import boto3
            from botocore.config import Config as BotoConfig
            _s3_client = boto3.client(
                service_name='s3',
                endpoint_url=S3_ENDPOINT,
                region_name=S3_REGION,
                aws_access_key_id=S3_ACCESS_KEY,
                aws_secret_access_key=S3_SECRET_KEY,
                config=BotoConfig(s3={'addressing_style': 'path'} if S3_FORCE_PATH else {})
            )
            logger.info("S3 клиент инициализирован")
        except Exception as e:
            logger.warning(f"S3 недоступен: {e}")
    return _s3_client
 
 
def upload_to_s3_background(file_data: bytes, table_name: str, message_id: int, filename: str, media_type: str):
    """Загружает файл в S3 в фоновом потоке (не блокирует бот)."""
    try:
        s3 = get_s3_client()
        if not s3:
            return
 
        import hashlib
        file_hash = hashlib.md5(file_data).hexdigest()[:12]
        safe_name = re.sub(r'[^\w\-.]', '_', filename) if filename else f"{media_type}_{file_hash}"
        s3_key = f"tg/{table_name}/{message_id}_{safe_name}"
 
        # Определяем content_type
        content_types = {
            'photo': 'image/jpeg', 'image': 'image/jpeg', 'pdf': 'application/pdf',
            'excel': 'application/vnd.openxmlformats-officedocument.spreadsheetml.sheet',
            'word': 'application/vnd.openxmlformats-officedocument.wordprocessingml.document',
            'powerpoint': 'application/vnd.openxmlformats-officedocument.presentationml.presentation',
            'video': 'video/mp4', 'voice': 'audio/ogg', 'audio': 'audio/mpeg',
        }
        content_type = content_types.get(media_type, 'application/octet-stream')
 
        s3.put_object(
            Bucket=S3_BUCKET,
            Key=s3_key,
            Body=file_data,
            ContentType=content_type
        )
 
        # Обновляем storage_path в БД
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute(
                    f"UPDATE {table_name} SET storage_path = %s WHERE message_id = %s",
                    (f"s3://{S3_BUCKET}/{s3_key}", message_id)
                )
                conn.commit()
        finally:
            conn.close()
 
        logger.info(f"S3 upload: {s3_key} ({len(file_data)} bytes)")
 
    except Exception as e:
        logger.warning(f"S3 upload failed (non-critical): {e}")
 

# ============================================================
# РАБОТА С БД
# ============================================================

def get_db_connection():
    """Создаёт подключение к PostgreSQL."""
    return psycopg2.connect(
        host=DB_HOST,
        port=DB_PORT,
        dbname=DB_NAME,
        user=DB_USER,
        password=DB_PASSWORD
    )


def sanitize_table_name(chat_id: int, chat_title: str) -> str:
    """Создаёт безопасное имя таблицы из ID и названия чата."""
    translit_map = {
        'а': 'a', 'б': 'b', 'в': 'v', 'г': 'g', 'д': 'd', 'е': 'e', 'ё': 'e',
        'ж': 'zh', 'з': 'z', 'и': 'i', 'й': 'y', 'к': 'k', 'л': 'l', 'м': 'm',
        'н': 'n', 'о': 'o', 'п': 'p', 'р': 'r', 'с': 's', 'т': 't', 'у': 'u',
        'ф': 'f', 'х': 'h', 'ц': 'ts', 'ч': 'ch', 'ш': 'sh', 'щ': 'sch',
        'ъ': '', 'ы': 'y', 'ь': '', 'э': 'e', 'ю': 'yu', 'я': 'ya'
    }
    
    title_lower = chat_title.lower()
    transliterated = ''.join(translit_map.get(c, c) for c in title_lower)
    safe_title = re.sub(r'[^a-z0-9]+', '_', transliterated)
    safe_title = re.sub(r'_+', '_', safe_title).strip('_')
    safe_title = safe_title[:30] if safe_title else "unnamed"
    
    return f"tg_chat_{abs(chat_id)}_{safe_title}"


def ensure_table_exists(chat_id: int, chat_title: str) -> str:
    """Создаёт таблицу для чата, если она не существует."""
    table_name = sanitize_table_name(chat_id, chat_title)
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(sql.SQL("""
                CREATE TABLE IF NOT EXISTS {} (
                    id SERIAL PRIMARY KEY,
                    message_id BIGINT NOT NULL,
                    user_id BIGINT,
                    username VARCHAR(255),
                    first_name VARCHAR(255),
                    last_name VARCHAR(255),
                    message_text TEXT,
                    message_type VARCHAR(50) DEFAULT 'text',
                    reply_to_message_id BIGINT,
                    forward_from_user_id BIGINT,
                    media_file_id TEXT,
                    media_analysis TEXT,
                    content_text TEXT,
                    timestamp TIMESTAMPTZ NOT NULL,
                    created_at TIMESTAMPTZ DEFAULT NOW(),
                    UNIQUE(message_id)
                )
            """).format(sql.Identifier(table_name)))
            
            cur.execute(sql.SQL("""
                CREATE INDEX IF NOT EXISTS {} ON {} (timestamp)
            """).format(
                sql.Identifier(f"idx_{table_name}_timestamp"),
                sql.Identifier(table_name)
            ))
            
            cur.execute(sql.SQL("""
                CREATE INDEX IF NOT EXISTS {} ON {} (user_id)
            """).format(
                sql.Identifier(f"idx_{table_name}_user_id"),
                sql.Identifier(table_name)
            ))
            
            cur.execute(sql.SQL("""
                CREATE INDEX IF NOT EXISTS {} ON {} USING gin(to_tsvector('russian', COALESCE(message_text, '') || ' ' || COALESCE(media_analysis, '') || ' ' || COALESCE(content_text, '')))
            """).format(
                sql.Identifier(f"idx_{table_name}_fts"),
                sql.Identifier(table_name)
            ))
            
            # Таблица метаданных чатов
            cur.execute("""
                CREATE TABLE IF NOT EXISTS tg_chats_metadata (
                    chat_id BIGINT PRIMARY KEY,
                    chat_title VARCHAR(255),
                    table_name VARCHAR(100),
                    chat_type VARCHAR(50),
                    added_at TIMESTAMPTZ DEFAULT NOW(),
                    last_message_at TIMESTAMPTZ,
                    total_messages INTEGER DEFAULT 0
                )
            """)
            
            # Таблица ролей пользователей
            cur.execute("""
                CREATE TABLE IF NOT EXISTS tg_user_roles (
                    id SERIAL PRIMARY KEY,
                    user_id BIGINT NOT NULL,
                    chat_id BIGINT NOT NULL,
                    username VARCHAR(255),
                    first_name VARCHAR(255),
                    last_name VARCHAR(255),
                    role VARCHAR(255),
                    created_at TIMESTAMPTZ DEFAULT NOW(),
                    updated_at TIMESTAMPTZ DEFAULT NOW(),
                    UNIQUE(user_id, chat_id)
                )
            """)
            
            cur.execute("""
                CREATE INDEX IF NOT EXISTS idx_user_roles_user ON tg_user_roles(user_id)
            """)
            
            cur.execute("""
                CREATE INDEX IF NOT EXISTS idx_user_roles_chat ON tg_user_roles(chat_id)
            """)
            
            cur.execute("""
                INSERT INTO tg_chats_metadata (chat_id, chat_title, table_name, chat_type, last_message_at)
                VALUES (%s, %s, %s, %s, NOW())
                ON CONFLICT (chat_id) DO UPDATE SET
                    chat_title = EXCLUDED.chat_title,
                    last_message_at = NOW()
            """, (chat_id, chat_title, table_name, "group"))
            
            conn.commit()
            logger.info(f"Таблица {table_name} готова для чата '{chat_title}'")
            
    finally:
        conn.close()
    
    return table_name


def save_message(table_name: str, message_data: dict):
    """Сохраняет сообщение в таблицу чата."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(sql.SQL("""
                INSERT INTO {} (
                    message_id, user_id, username, first_name, last_name,
                    message_text, message_type, reply_to_message_id,
                    forward_from_user_id, media_file_id, media_analysis, content_text, timestamp
                ) VALUES (
                    %(message_id)s, %(user_id)s, %(username)s, %(first_name)s, %(last_name)s,
                    %(message_text)s, %(message_type)s, %(reply_to_message_id)s,
                    %(forward_from_user_id)s, %(media_file_id)s, %(media_analysis)s, %(content_text)s, %(timestamp)s
                )
                ON CONFLICT (message_id) DO UPDATE SET
                    message_text = EXCLUDED.message_text,
                    media_analysis = EXCLUDED.media_analysis,
                    content_text = EXCLUDED.content_text
            """).format(sql.Identifier(table_name)), message_data)
            conn.commit()
            # Canonical zone
            try:
                from canonical_helper import insert_source_document_tg
                chat_meta_cur = conn.cursor()
                chat_meta_cur.execute(
                    "SELECT chat_title FROM tg_chats_metadata WHERE table_name = %s",
                    (table_name,)
                )
                row = chat_meta_cur.fetchone()
                chat_title = row[0] if row else table_name
                chat_meta_cur.close()
                
                with conn.cursor() as canon_cur:
                    insert_source_document_tg(canon_cur, table_name, chat_title, message_data)
                    conn.commit()
            except Exception as e:
                logger.warning(f"Canonical insert error: {e}")
    finally:
        conn.close()


# ============================================================
# РАБОТА С РОЛЯМИ ПОЛЬЗОВАТЕЛЕЙ
# ============================================================

def get_user_role(user_id: int, chat_id: int) -> str | None:
    """Получает роль пользователя в чате."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT role FROM tg_user_roles
                WHERE user_id = %s AND chat_id = %s
            """, (user_id, chat_id))
            result = cur.fetchone()
            return result[0] if result else None
    except Exception as e:
        logger.error(f"Ошибка получения роли: {e}")
        return None
    finally:
        conn.close()


def set_user_role(user_id: int, chat_id: int, role: str, username: str = None, first_name: str = None, last_name: str = None):
    """Устанавливает роль пользователя в чате."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO tg_user_roles (user_id, chat_id, username, first_name, last_name, role, updated_at)
                VALUES (%s, %s, %s, %s, %s, %s, NOW())
                ON CONFLICT (user_id, chat_id) DO UPDATE SET
                    role = EXCLUDED.role,
                    username = COALESCE(EXCLUDED.username, tg_user_roles.username),
                    first_name = COALESCE(EXCLUDED.first_name, tg_user_roles.first_name),
                    last_name = COALESCE(EXCLUDED.last_name, tg_user_roles.last_name),
                    updated_at = NOW()
            """, (user_id, chat_id, username, first_name, last_name, role))
            conn.commit()
            logger.info(f"Роль '{role}' установлена для пользователя {user_id} в чате {chat_id}")
    except Exception as e:
        logger.error(f"Ошибка установки роли: {e}")
    finally:
        conn.close()

# ============================================================
# НАСТРОЙКИ РАССЫЛКИ ПОЛНОГО АНАЛИЗА В ЛИЧКУ
# ============================================================

def get_user_analysis_setting(user_id: int) -> bool:
    """Проверяет, включена ли у пользователя рассылка полного анализа."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT send_full_analysis 
                FROM tg_full_analysis_settings 
                WHERE user_id = %s
            """, (user_id,))
            row = cur.fetchone()
            return row[0] if row else False
    except Exception as e:
        logger.error(f"Ошибка получения настройки анализа: {e}")
        return False
    finally:
        conn.close()


def set_user_analysis_setting(user_id: int, username: str, first_name: str, enabled: bool):
    """Устанавливает настройку рассылки полного анализа."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO tg_full_analysis_settings 
                    (user_id, username, first_name, send_full_analysis, updated_at)
                VALUES (%s, %s, %s, %s, NOW())
                ON CONFLICT (user_id) DO UPDATE SET
                    send_full_analysis = EXCLUDED.send_full_analysis,
                    username = COALESCE(EXCLUDED.username, tg_full_analysis_settings.username),
                    first_name = COALESCE(EXCLUDED.first_name, tg_full_analysis_settings.first_name),
                    updated_at = NOW()
            """, (user_id, username, first_name, enabled))
            conn.commit()
    except Exception as e:
        logger.error(f"Ошибка установки настройки анализа: {e}")
    finally:
        conn.close()


def get_users_with_full_analysis_enabled() -> list:
    """Возвращает список user_id с включённой рассылкой."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT user_id, username, first_name
                FROM tg_full_analysis_settings 
                WHERE send_full_analysis = TRUE
            """)
            return cur.fetchall()
    except Exception as e:
        logger.error(f"Ошибка получения списка пользователей: {e}")
        return []
    finally:
        conn.close()


def get_users_without_roles(chat_id: int, table_name: str) -> list:
    """Получает список пользователей без ролей в чате."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(sql.SQL("""
                SELECT DISTINCT m.user_id, m.username, m.first_name, m.last_name
                FROM {} m
                LEFT JOIN tg_user_roles r ON m.user_id = r.user_id AND r.chat_id = %s
                WHERE m.user_id IS NOT NULL 
                AND r.role IS NULL
                AND m.timestamp > NOW() - INTERVAL '30 days'
                ORDER BY m.first_name
            """).format(sql.Identifier(table_name)), (chat_id,))
            return cur.fetchall()
    except Exception as e:
        logger.error(f"Ошибка получения пользователей без ролей: {e}")
        return []
    finally:
        conn.close()


# ============================================================
# КОНТЕКСТ ЧАТА
# ============================================================

def get_full_chat_context(table_name: str, chat_id: int, chat_title: str, hours: int = 192) -> str:
    """Получает полный контекст чата с ролями и связанными сообщениями за 8 дней."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            # Получаем участников с ролями
            cur.execute(sql.SQL("""
                SELECT DISTINCT 
                    m.user_id,
                    COALESCE(m.first_name, m.username, 'Неизвестный') as name,
                    m.last_name,
                    r.role
                FROM {} m
                LEFT JOIN tg_user_roles r ON m.user_id = r.user_id AND r.chat_id = %s
                WHERE m.timestamp > NOW() - INTERVAL '%s hours'
                AND m.user_id IS NOT NULL
            """).format(sql.Identifier(table_name)), (chat_id, hours))
            
            participants = cur.fetchall()
            
            # Получаем сообщения за период
            cur.execute(sql.SQL("""
                SELECT 
                    m.message_id,
                    m.user_id,
                    COALESCE(m.first_name, m.username, 'Неизвестный') as first_name,
                    m.last_name,
                    r.role,
                    m.message_text,
                    m.media_analysis,
                    m.message_type,
                    m.reply_to_message_id,
                    m.timestamp
                FROM {} m
                LEFT JOIN tg_user_roles r ON m.user_id = r.user_id AND r.chat_id = %s
                WHERE m.timestamp > NOW() - INTERVAL '%s hours'
                AND (m.message_text IS NOT NULL AND m.message_text != '' 
                     OR m.media_analysis IS NOT NULL AND m.media_analysis != '')
                ORDER BY m.timestamp ASC
            """).format(sql.Identifier(table_name)), (chat_id, hours))
            
            messages = cur.fetchall()
            
            # Собираем ID для поиска связанных сообщений
            reply_ids = [m[8] for m in messages if m[8] is not None]
            message_ids = [m[0] for m in messages]
            
            # Находим связанные сообщения за пределами периода
            missing_ids = [rid for rid in reply_ids if rid not in message_ids]
            linked_messages = {}
            
            if missing_ids:
                placeholders = ','.join(['%s'] * len(missing_ids))
                cur.execute(sql.SQL(f"""
                    SELECT 
                        m.message_id,
                        COALESCE(m.first_name, m.username, 'Неизвестный') as first_name,
                        m.last_name,
                        r.role,
                        m.message_text,
                        m.media_analysis,
                        m.message_type,
                        m.timestamp
                    FROM {{}} m
                    LEFT JOIN tg_user_roles r ON m.user_id = r.user_id AND r.chat_id = %s
                    WHERE m.message_id IN ({placeholders})
                """).format(sql.Identifier(table_name)), [chat_id] + missing_ids)
                
                for row in cur.fetchall():
                    linked_messages[row[0]] = row
            
            if not messages and not participants:
                return ""
            
            # Формируем контекст
            context_parts = []
            
            # Информация о чате
            context_parts.append(f"=== ЧАТ: {chat_title} ===\n")
            
            # Участники с ролями
            context_parts.append("УЧАСТНИКИ ЧАТА:")
            for user_id, name, last_name, role in participants:
                full_name = f"{name} {last_name}" if last_name else name
                role_str = f" — {role}" if role else " — роль не указана"
                context_parts.append(f"  • {full_name}{role_str}")
            context_parts.append("")
            
            # Сообщения
            context_parts.append("=== ИСТОРИЯ СООБЩЕНИЙ (последние 8 дней) ===\n")
            
            for msg_id, user_id, first_name, last_name, role, text, analysis, msg_type, reply_to, ts in messages:
                date_str = ts.strftime("%d.%m.%Y")
                time_str = ts.strftime("%H:%M")
                full_name = f"{first_name} {last_name}" if last_name else first_name
                role_str = f" [{role}]" if role else ""
                
                msg_parts = [f"[{date_str} {time_str}] {full_name}{role_str}:"]
                
                # Если это ответ на другое сообщение
                if reply_to:
                    linked = None
                    # Сначала ищем в связанных сообщениях за пределами периода
                    if reply_to in linked_messages:
                        linked = linked_messages[reply_to]
                        linked_name = f"{linked[1]} {linked[2]}" if linked[2] else linked[1]
                        linked_role = f" [{linked[3]}]" if linked[3] else ""
                        linked_date = linked[7].strftime("%d.%m.%Y %H:%M")
                        linked_content = linked[4] if linked[4] else linked[5] if linked[5] else "[медиа]"
                        linked_content = linked_content[:300] + "..." if len(linked_content) > 300 else linked_content
                    else:
                        # Ищем в текущих сообщениях
                        for m in messages:
                            if m[0] == reply_to:
                                linked_name = f"{m[2]} {m[3]}" if m[3] else m[2]
                                linked_role = f" [{m[4]}]" if m[4] else ""
                                linked_date = m[9].strftime("%d.%m.%Y %H:%M")
                                linked_content = m[5] if m[5] else m[6] if m[6] else "[медиа]"
                                linked_content = linked_content[:300] + "..." if len(linked_content) > 300 else linked_content
                                linked = True
                                break
                    
                    if linked:
                        msg_parts.append(f"  ↳ В ОТВЕТ НА ({linked_date}, {linked_name}{linked_role}):")
                        msg_parts.append(f"    \"{linked_content}\"")
                
                if text and text.strip():
                    msg_parts.append(f"  {text[:3000]}")
                
                if analysis and analysis.strip():
                    analysis_short = analysis[:1600] + "..." if len(analysis) > 1600 else analysis
                    msg_parts.append(f"  [АНАЛИЗ {msg_type.upper()}]: {analysis_short}")
                
                context_parts.append("\n".join(msg_parts))
                context_parts.append("")
            
            return "\n".join(context_parts)
            
    except Exception as e:
        logger.error(f"Ошибка получения контекста чата: {e}")
        return ""
    finally:
        conn.close()


# ============================================================
# ПОСТРОЕНИЕ ПРОМПТА
# ============================================================

def build_analysis_prompt(doc_type: str, doc_content: str, context: str, filename: str = "") -> str:
    """Создаёт промпт для анализа документа с учётом контекста чата и знаний о компании."""
    
    company_profile = get_company_profile()
    
    prompt = f"""{company_profile}

Ты участник рабочего чата компании Фрумелад. Коллега отправил документ — дай краткий анализ.

ПРАВИЛА ОТВЕТА:
- Пиши сплошным текстом, как сообщение в чате коллеге.
- НИКОГДА не используй заголовки (##, **Заголовок**), списки, буллеты (-, •, *), нумерованные пункты, разделы, разделители (---).
- Будь краток пропорционально сложности: простое фото или скриншот — 3-5 предложений, накладная или счёт — 5-7 предложений, сложный многостраничный документ (договор, спецификация, прайс) — до 15 предложений. Никогда не превышай 15 предложений.
- Начни с главного: что это и зачем это участникам чата.
- Упомяни ключевые цифры, даты, контрагентов, суммы — если они есть.
- Не повторяй одну мысль разными словами.
- Если нужно действие — скажи кратко в конце.

"""
    
    if context:
        prompt += f"""Контекст обсуждения в чате:
{context}

Анализируй документ ИМЕННО в контексте этого обсуждения. Отвечай на тот вопрос, который обсуждался.

"""
    
    prompt += f"""Тип: {doc_type}
Файл: {filename}
{doc_content if doc_content != "[Изображение прикреплено]" else ""}

Краткий анализ:"""
    
    return prompt


# ============================================================
# ИЗВЛЕЧЕНИЕ ТЕКСТОВОГО СОДЕРЖИМОГО
# ============================================================

async def extract_text_from_image(image_data: bytes, media_type: str) -> str:
    """Извлекает текст из изображения с помощью OCR через Claude Vision."""
    if not gpt_client:
        return ""

    try:
        base64_image = base64.standard_b64encode(image_data).decode("utf-8")

        # Используем Claude Vision для OCR
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4096,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{media_type};base64,{base64_image}"
                            }
                        },
                        {
                            "type": "text",
                            "text": "Извлеки весь текст, который видишь на этом изображении. Верни только текст, без дополнительных комментариев. Если текста нет, верни пустую строку."
                        }
                    ],
                }
            ],
        )

        extracted_text = response.choices[0].message.content.strip()
        logger.info(f"Текст извлечен из изображения: {len(extracted_text)} символов")
        return extracted_text

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из изображения: {e}")
        return ""


async def extract_text_from_pdf(pdf_data: bytes) -> str:
    """Извлекает текст из PDF файла."""
    try:
        # Пробуем извлечь текст напрямую
        try:
            import PyPDF2
            import io

            pdf_reader = PyPDF2.PdfReader(io.BytesIO(pdf_data))
            text_parts = []

            for page_num, page in enumerate(pdf_reader.pages[:50], 1):  # Максимум 50 страниц
                page_text = page.extract_text()
                if page_text.strip():
                    text_parts.append(f"=== Страница {page_num} ===\n{page_text}")

            if text_parts:
                return "\n\n".join(text_parts)
        except Exception as e:
            logger.warning(f"PyPDF2 не смог извлечь текст: {e}")

        # Если PyPDF2 не сработал, используем Claude Vision через pdf2image
        try:
            from pdf2image import convert_from_bytes
            images = convert_from_bytes(pdf_data, first_page=1, last_page=20)

            all_text = []
            for i, image in enumerate(images, 1):
                import io
                img_byte_arr = io.BytesIO()
                image.save(img_byte_arr, format='PNG')
                img_bytes = img_byte_arr.getvalue()

                page_text = await extract_text_from_image(img_bytes, "image/png")
                if page_text:
                    all_text.append(f"=== Страница {i} ===\n{page_text}")

            return "\n\n".join(all_text)
        except Exception as e:
            logger.error(f"Не удалось извлечь текст через OCR: {e}")
            return ""

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из PDF: {e}")
        return ""


async def extract_text_from_word(file_data: bytes) -> str:
    """Извлекает текст из Word документа."""
    try:
        import io
        from docx import Document

        doc = Document(io.BytesIO(file_data))

        text_parts = []

        # Извлекаем параграфы
        for para in doc.paragraphs:
            if para.text.strip():
                text_parts.append(para.text)

        # Извлекаем таблицы
        for table in doc.tables:
            table_text = []
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells)
                if row_text.strip():
                    table_text.append(row_text)
            if table_text:
                text_parts.append("\n" + "\n".join(table_text))

        return "\n".join(text_parts)

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из Word: {e}")
        return ""


async def extract_csv_from_excel(file_data: bytes, filename: str = "") -> str:
    """Извлекает данные из Excel в формате CSV."""
    try:
        import io
        is_xls = filename.lower().endswith('.xls') and not filename.lower().endswith('.xlsx')

        csv_parts = []

        if is_xls:
            # Старый формат .xls
            try:
                import xlrd
                wb = xlrd.open_workbook(file_contents=file_data)

                for sheet_name in wb.sheet_names():
                    sheet = wb.sheet_by_name(sheet_name)
                    csv_parts.append(f"=== {sheet_name} ===")

                    for row_idx in range(sheet.nrows):
                        row_values = []
                        for col_idx in range(sheet.ncols):
                            cell_value = sheet.cell_value(row_idx, col_idx)
                            # Экранируем значения для CSV
                            if isinstance(cell_value, str) and (',' in cell_value or '"' in cell_value or '\n' in cell_value):
                                cell_value = f'"{cell_value.replace(chr(34), chr(34)+chr(34))}"'
                            row_values.append(str(cell_value) if cell_value else "")
                        csv_parts.append(",".join(row_values))
                    csv_parts.append("")  # Пустая строка между листами

            except Exception as e:
                logger.error(f"Ошибка чтения .xls: {e}")
                return ""
        else:
            # Новый формат .xlsx
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)

                for sheet_name in wb.sheetnames:
                    sheet = wb[sheet_name]
                    csv_parts.append(f"=== {sheet_name} ===")

                    for row in sheet.iter_rows(values_only=True):
                        row_values = []
                        for cell in row:
                            # Экранируем значения для CSV
                            cell_str = str(cell) if cell is not None else ""
                            if ',' in cell_str or '"' in cell_str or '\n' in cell_str:
                                cell_str = f'"{cell_str.replace(chr(34), chr(34)+chr(34))}"'
                            row_values.append(cell_str)
                        csv_parts.append(",".join(row_values))
                    csv_parts.append("")  # Пустая строка между листами

                wb.close()
            except Exception as e:
                logger.warning(f"openpyxl не смог открыть, пробуем xlrd: {e}")
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_contents=file_data)

                    for sheet_name in wb.sheet_names():
                        sheet = wb.sheet_by_name(sheet_name)
                        csv_parts.append(f"=== {sheet_name} ===")

                        for row_idx in range(sheet.nrows):
                            row_values = []
                            for col_idx in range(sheet.ncols):
                                cell_value = sheet.cell_value(row_idx, col_idx)
                                if isinstance(cell_value, str) and (',' in cell_value or '"' in cell_value or '\n' in cell_value):
                                    cell_value = f'"{cell_value.replace(chr(34), chr(34)+chr(34))}"'
                                row_values.append(str(cell_value) if cell_value else "")
                            csv_parts.append(",".join(row_values))
                        csv_parts.append("")
                except Exception as e2:
                    logger.error(f"Ошибка чтения Excel обоими методами: {e2}")
                    return ""

        return "\n".join(csv_parts)

    except Exception as e:
        logger.error(f"Ошибка извлечения CSV из Excel: {e}")
        return ""


async def extract_text_from_pptx(file_data: bytes) -> str:
    """Извлекает текст из PowerPoint презентации."""
    try:
        import io
        from pptx import Presentation

        prs = Presentation(io.BytesIO(file_data))

        text_parts = []

        for slide_num, slide in enumerate(prs.slides, 1):
            slide_text = [f"=== Слайд {slide_num} ==="]

            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_text.append(shape.text)

                # Извлекаем текст из таблиц
                if hasattr(shape, "table"):
                    table = shape.table
                    for row in table.rows:
                        row_text = " | ".join(cell.text.strip() for cell in row.cells)
                        if row_text:
                            slide_text.append(row_text)

            if len(slide_text) > 1:  # Если есть текст помимо заголовка
                text_parts.append("\n".join(slide_text))

        return "\n\n".join(text_parts)

    except Exception as e:
        logger.error(f"Ошибка извлечения текста из PowerPoint: {e}")
        return ""


async def extract_transcript_from_audio(audio_path: str) -> str:
    """Извлекает транскрипт из аудио файла с помощью Whisper."""
    try:
        import whisper

        model = whisper.load_model("base")
        result = model.transcribe(audio_path, language="ru", fp16=False)

        transcript = result.get("text", "").strip()
        logger.info(f"Аудио транскрибировано: {len(transcript)} символов")
        return transcript

    except Exception as e:
        logger.error(f"Ошибка транскрибирования аудио: {e}")
        return ""


# ============================================================
# АНАЛИЗ ДОКУМЕНТОВ
# ============================================================

async def analyze_image_with_gpt(image_data: bytes, media_type: str, context: str = "", filename: str = "") -> str:
    """Анализирует изображение через Claude Vision."""
    if not gpt_client:
        return ""
    
    try:
        base64_image = base64.standard_b64encode(image_data).decode("utf-8")
        
        prompt = build_analysis_prompt("Изображение", "[Изображение прикреплено]", context, filename)
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=1500,
            messages=[
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{media_type};base64,{base64_image}"
                            }
                        },
                        {
                            "type": "text",
                            "text": prompt
                        }
                    ],
                }
            ],
        )
        
        analysis = response.choices[0].message.content
        logger.info(f"Изображение проанализировано: {len(analysis)} символов")
        return analysis
        
    except Exception as e:
        logger.error(f"Ошибка анализа изображения: {e}")
        return ""


async def analyze_pdf_with_gpt(pdf_data: bytes, filename: str = "", context: str = "") -> str:
    """Анализирует PDF — все страницы одним запросом."""
    if not gpt_client:
        return ""
    
    try:
        try:
            from pdf2image import convert_from_bytes
            images = convert_from_bytes(pdf_data, first_page=1, last_page=10)
        except Exception as e:
            logger.warning(f"Не удалось конвертировать PDF в изображения: {e}")
            base64_pdf = base64.standard_b64encode(pdf_data).decode("utf-8")
            
            prompt = build_analysis_prompt("PDF документ", "[PDF документ прикреплён]", context, filename)
            
            response = gpt_client.chat.completions.create(
                model="openai/gpt-4.1",
                max_tokens=4500,
                messages=[
                    {
                        "role": "user",
                        "content": [
                            {
                                "type": "document",
                                "source": {
                                    "type": "base64",
                                    "media_type": "application/pdf",
                                    "data": base64_pdf,
                                },
                            },
                            {
                                "type": "text",
                                "text": prompt
                            }
                        ],
                    }
                ],
            )
            return response.choices[0].message.content
        
        # Собираем все страницы в один запрос
        import io
        content_parts = []
        
        for i, image in enumerate(images):
            img_byte_arr = io.BytesIO()
            image.save(img_byte_arr, format='PNG')
            img_bytes = img_byte_arr.getvalue()
            base64_image = base64.standard_b64encode(img_bytes).decode("utf-8")
            
            content_parts.append({
                "type": "image_url",
                "image_url": {
                    "url": f"data:image/png;base64,{base64_image}"
                }
            })
        
        prompt = build_analysis_prompt(
            "PDF документ",
            f"[PDF документ из {len(images)} страниц прикреплён]",
            context,
            filename
        )
        prompt += f"\n\nДокумент содержит {len(images)} страниц. Проанализируй весь документ целиком, учитывая содержимое всех страниц."
        
        content_parts.append({
            "type": "text",
            "text": prompt
        })
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4500,
            messages=[{"role": "user", "content": content_parts}],
        )
        
        logger.info(f"PDF проанализирован целиком ({len(images)} страниц)")
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Ошибка анализа PDF: {e}")
        return ""


async def analyze_excel_with_gpt(file_data: bytes, filename: str = "", context: str = "") -> str:
    """Анализирует Excel файл через Claude. Поддерживает .xlsx и .xls форматы."""
    if not gpt_client:
        return ""
    
    try:
        import io
        all_text = []
        
        # Определяем формат по расширению или пробуем оба
        is_xls = filename.lower().endswith('.xls') and not filename.lower().endswith('.xlsx')
        
        if is_xls:
            # Старый формат .xls
            try:
                import xlrd
                wb = xlrd.open_workbook(file_contents=file_data)
                
                for sheet_name in wb.sheet_names()[:5]:
                    sheet = wb.sheet_by_name(sheet_name)
                    all_text.append(f"=== Лист: {sheet_name} ===")
                    
                    rows_count = 0
                    for row_idx in range(min(sheet.nrows, 200)):
                        row_values = [str(sheet.cell_value(row_idx, col_idx)) if sheet.cell_value(row_idx, col_idx) else "" for col_idx in range(sheet.ncols)]
                        if any(row_values):
                            all_text.append(" | ".join(row_values))
                            rows_count += 1
                    
                    if rows_count == 200:
                        all_text.append("... (данные обрезаны)")
                        
            except Exception as e:
                logger.error(f"Ошибка чтения .xls: {e}")
                return ""
        else:
            # Новый формат .xlsx
            try:
                from openpyxl import load_workbook
                wb = load_workbook(io.BytesIO(file_data), read_only=True, data_only=True)
                
                for sheet_name in wb.sheetnames[:5]:
                    sheet = wb[sheet_name]
                    all_text.append(f"=== Лист: {sheet_name} ===")
                    
                    rows_count = 0
                    for row in sheet.iter_rows(max_row=200, values_only=True):
                        row_values = [str(cell) if cell is not None else "" for cell in row]
                        if any(row_values):
                            all_text.append(" | ".join(row_values))
                            rows_count += 1
                    
                    if rows_count == 200:
                        all_text.append("... (данные обрезаны)")
                
                wb.close()
            except Exception as e:
                # Может это .xls с неправильным расширением - пробуем xlrd
                logger.warning(f"openpyxl не смог открыть, пробуем xlrd: {e}")
                try:
                    import xlrd
                    wb = xlrd.open_workbook(file_contents=file_data)
                    
                    for sheet_name in wb.sheet_names()[:5]:
                        sheet = wb.sheet_by_name(sheet_name)
                        all_text.append(f"=== Лист: {sheet_name} ===")
                        
                        for row_idx in range(min(sheet.nrows, 200)):
                            row_values = [str(sheet.cell_value(row_idx, col_idx)) if sheet.cell_value(row_idx, col_idx) else "" for col_idx in range(sheet.ncols)]
                            if any(row_values):
                                all_text.append(" | ".join(row_values))
                except Exception as e2:
                    logger.error(f"Ошибка чтения Excel обоими методами: {e2}")
                    return ""
        
        excel_content = "\n".join(all_text)
        
        if len(excel_content) > 15000:
            excel_content = excel_content[:15000] + "\n... (данные обрезаны)"
        
        if not excel_content.strip():
            return "Файл пуст или не удалось прочитать содержимое."
        
        prompt = build_analysis_prompt("Excel таблица", f"Содержимое файла:\n{excel_content}", context, filename)
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4500,
            messages=[{"role": "user", "content": prompt}],
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Ошибка анализа Excel: {e}")
        return ""


async def analyze_word_with_gpt(file_data: bytes, filename: str = "", context: str = "") -> str:
    """Анализирует Word файл через Claude."""
    if not gpt_client:
        return ""
    
    try:
        import io
        from docx import Document
        
        doc = Document(io.BytesIO(file_data))
        
        paragraphs = []
        for para in doc.paragraphs[:500]:
            if para.text.strip():
                paragraphs.append(para.text)
        
        for table in doc.tables[:10]:
            for row in table.rows:
                row_text = " | ".join(cell.text.strip() for cell in row.cells if cell.text.strip())
                if row_text:
                    paragraphs.append(row_text)
        
        word_content = "\n".join(paragraphs)
        
        if len(word_content) > 15000:
            word_content = word_content[:15000] + "\n... (текст обрезан)"
        
        prompt = build_analysis_prompt("Word документ", f"Содержимое документа:\n{word_content}", context, filename)
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4500,
            messages=[{"role": "user", "content": prompt}],
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Ошибка анализа Word: {e}")
        return ""


async def analyze_pptx_with_gpt(file_data: bytes, filename: str = "", context: str = "") -> str:
    """Анализирует PowerPoint файл через Claude."""
    if not gpt_client:
        return ""
    
    try:
        import io
        from pptx import Presentation
        
        prs = Presentation(io.BytesIO(file_data))
        
        slides_text = []
        for i, slide in enumerate(prs.slides[:30], 1):
            slide_content = [f"=== Слайд {i} ==="]
            
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if len(slide_content) > 1:
                slides_text.append("\n".join(slide_content))
        
        pptx_content = "\n\n".join(slides_text)
        
        if len(pptx_content) > 15000:
            pptx_content = pptx_content[:15000] + "\n... (текст обрезан)"
        
        prompt = build_analysis_prompt("PowerPoint презентация", f"Содержимое презентации:\n{pptx_content}", context, filename)
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4500,
            messages=[{"role": "user", "content": prompt}],
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Ошибка анализа PowerPoint: {e}")
        return ""

async def analyze_video_with_gemini(file_data: bytes, filename: str = "", context: str = "") -> str:
    """Анализирует видео через Gemini 3 Flash (поддерживает видео напрямую)."""
    import requests
    import base64
    
    ROUTERAI_API_KEY = os.getenv("ROUTERAI_API_KEY")
    ROUTERAI_BASE_URL = os.getenv("ROUTERAI_BASE_URL", "https://routerai.ru/api/v1")
    
    if not ROUTERAI_API_KEY:
        logger.warning("ROUTERAI_API_KEY не установлен — анализ видео через Gemini недоступен")
        # Fallback на старый метод с Whisper + Claude
        return await analyze_video_with_whisper(file_data, filename, context)
    
    try:
        # Кодируем видео в base64
        video_base64 = base64.standard_b64encode(file_data).decode("utf-8")
        
        # Определяем mime type
        ext = filename.lower().split('.')[-1] if filename else 'mp4'
        mime_types = {
            'mp4': 'video/mp4',
            'avi': 'video/x-msvideo',
            'mov': 'video/quicktime',
            'mkv': 'video/x-matroska',
            'webm': 'video/webm'
        }
        mime_type = mime_types.get(ext, 'video/mp4')
        
        # Строим промпт
        prompt = build_analysis_prompt("Видео", "[Видео прикреплено]", context, filename)
        
        # Запрос к Gemini через RouterAI
        url = f"{ROUTERAI_BASE_URL}/chat/completions"
        headers = {
            "Authorization": f"Bearer {ROUTERAI_API_KEY}",
            "Content-Type": "application/json"
        }
        
        data = {
            "model": "google/gemini-3-flash-preview",
            "messages": [
                {
                    "role": "user",
                    "content": [
                        {
                            "type": "text",
                            "text": prompt
                        },
                        {
                            "type": "image_url",
                            "image_url": {
                                "url": f"data:{mime_type};base64,{video_base64}"
                            }
                        }
                    ]
                }
            ],
            "max_tokens": 4000
        }
        
        response = requests.post(url, headers=headers, json=data, timeout=120)
        result = response.json()
        
        if "choices" in result and len(result["choices"]) > 0:
            analysis = result["choices"][0]["message"]["content"]
            logger.info(f"Видео проанализировано через Gemini: {len(analysis)} символов")
            return analysis
        else:
            logger.error(f"Ошибка Gemini API: {result}")
            return ""
            
    except Exception as e:
        logger.error(f"Ошибка анализа видео через Gemini: {e}")
        return ""


async def analyze_video_with_whisper(file_data: bytes, filename: str = "", context: str = "") -> str:
    """Fallback: анализирует видео через Whisper (только аудио) + Claude."""
    if not gpt_client:
        return ""
    
    try:
        import tempfile
        import subprocess
        import whisper
        
        with tempfile.NamedTemporaryFile(suffix='.mp4', delete=False) as video_file:
            video_file.write(file_data)
            video_path = video_file.name
        
        audio_path = video_path.replace('.mp4', '.wav')
        subprocess.run([
            'ffmpeg', '-i', video_path, '-vn', '-acodec', 'pcm_s16le',
            '-ar', '16000', '-ac', '1', audio_path, '-y'
        ], capture_output=True, timeout=120)
        
        transcript = ""
        if os.path.exists(audio_path) and os.path.getsize(audio_path) > 0:
            model = whisper.load_model("base")
            result = model.transcribe(audio_path, language="ru")
            transcript = result["text"]
        
        os.unlink(video_path)
        if os.path.exists(audio_path):
            os.unlink(audio_path)
        
        if not transcript.strip():
            return "Видео без речи или речь не распознана."
        
        if len(transcript) > 10000:
            transcript = transcript[:10000] + "... (транскрипция обрезана)"
        
        prompt = build_analysis_prompt("Видео (транскрипция аудио)", f"Транскрипция:\n{transcript}", context, filename)
        
        response = gpt_client.chat.completions.create(
            model="openai/gpt-4.1",
            max_tokens=4500,
            messages=[{"role": "user", "content": prompt}],
        )
        
        return response.choices[0].message.content
        
    except Exception as e:
        logger.error(f"Ошибка анализа видео через Whisper: {e}")
        return ""


# ============================================================
# ОБРАБОТКА МЕДИАФАЙЛОВ
# ============================================================

async def download_and_analyze_media(bot, message, table_name: str = None) -> tuple[str, str, str]:
    """Скачивает и анализирует медиафайл с учётом контекста чата.

    Возвращает: (media_type_str, media_analysis, content_text)
    """
    media_analysis = ""
    content_text = ""
    media_type_str = "media"

    try:
        file = None
        media_type = None
        filename = ""

        if message.photo:
            file = await bot.get_file(message.photo[-1].file_id)
            media_type = "image/jpeg"
            media_type_str = "photo"
            filename = "photo.jpg"
        elif message.voice:
            file = await bot.get_file(message.voice.file_id)
            media_type = "voice"
            media_type_str = "voice"
            filename = "voice.ogg"
        elif message.audio:
            file = await bot.get_file(message.audio.file_id)
            media_type = "audio"
            media_type_str = "audio"
            filename = message.audio.file_name or "audio.mp3"
        elif message.video:
            if message.video.file_size and message.video.file_size < 40 * 1024 * 1024:
                file = await bot.get_file(message.video.file_id)
                media_type = "video"
                media_type_str = "video"
                filename = "video.mp4"
            else:
                logger.warning("Видео слишком большое для анализа")
                return "video", "", ""
        elif message.document:
            doc = message.document
            filename = doc.file_name or ""
            filename_lower = filename.lower()

            if doc.mime_type and doc.mime_type.startswith("image/"):
                file = await bot.get_file(doc.file_id)
                media_type = doc.mime_type
                media_type_str = "image"
            elif doc.mime_type == "application/pdf" or filename_lower.endswith(".pdf"):
                file = await bot.get_file(doc.file_id)
                media_type = "application/pdf"
                media_type_str = "pdf"
            elif filename_lower.endswith(('.xlsx', '.xls')):
                file = await bot.get_file(doc.file_id)
                media_type = "excel"
                media_type_str = "excel"
            elif filename_lower.endswith(('.docx', '.doc')):
                file = await bot.get_file(doc.file_id)
                media_type = "word"
                media_type_str = "word"
            elif filename_lower.endswith(('.pptx', '.ppt')):
                file = await bot.get_file(doc.file_id)
                media_type = "powerpoint"
                media_type_str = "powerpoint"
            elif filename_lower.endswith(('.mp4', '.avi', '.mov', '.mkv')):
                if doc.file_size and doc.file_size < 40 * 1024 * 1024:
                    file = await bot.get_file(doc.file_id)
                    media_type = "video"
                    media_type_str = "video"
                else:
                    logger.warning("Видео слишком большое для анализа")
                    return "video", "", ""
            else:
                media_type_str = "document"
                return media_type_str, "", ""
        else:
            return media_type_str, "", ""

        if not file:
            return media_type_str, "", ""

        # Скачиваем файл
        file_data = await file.download_as_bytearray()
        
        # === S3 upload (фоновый, не блокирует анализ) ===
        if S3_BUCKET and len(file_data) > 0:
            threading.Thread(
                target=upload_to_s3_background,
                args=(bytes(file_data), table_name or "unknown", message.message_id, filename, media_type_str),
                daemon=True
            ).start()

        # Получаем полный контекст чата (8 дней = 192 часа)
        context = ""
        if table_name and message.chat:
            chat_context = get_full_chat_context(
                table_name,
                message.chat.id,
                message.chat.title or "Чат",
                192  # 8 дней
            )
            if chat_context:
                context = chat_context

        # Добавляем подпись если есть
        caption = message.caption or ""
        if caption:
            context += f"\n\n=== ПОДПИСЬ К ТЕКУЩЕМУ ФАЙЛУ ===\n{caption}"

        # Анализируем И извлекаем содержимое в зависимости от типа
        if media_type == "voice" or media_type == "audio":
            # Для голосовых сообщений и аудио - транскрибируем
            import tempfile
            suffix = '.ogg' if media_type == "voice" else '.mp3'
            with tempfile.NamedTemporaryFile(delete=False, suffix=suffix) as tmp:
                tmp.write(bytes(file_data))
                audio_path = tmp.name
            try:
                content_text = await extract_transcript_from_audio(audio_path)

                # Создаем анализ на основе транскрипта
                if content_text:
                    prompt = build_analysis_prompt(
                        "Голосовое сообщение" if media_type == "voice" else "Аудио файл",
                        f"Транскрипция:\n{content_text}",
                        context,
                        filename
                    )

                    if gpt_client:
                        response = gpt_client.chat.completions.create(
                            model="openai/gpt-4.1",
                            max_tokens=4500,
                            messages=[{"role": "user", "content": prompt}],
                        )
                        media_analysis = response.choices[0].message.content
                    else:
                        media_analysis = f"Транскрипция: {content_text}"
                else:
                    media_analysis = "Не удалось распознать речь в аудио."
            finally:
                if os.path.exists(audio_path):
                    os.unlink(audio_path)
        elif media_type == "application/pdf":
            media_analysis = await analyze_pdf_with_gpt(bytes(file_data), filename, context)
            content_text = await extract_text_from_pdf(bytes(file_data))
        elif media_type and media_type.startswith("image/"):
            media_analysis = await analyze_image_with_gpt(bytes(file_data), media_type, context, filename)
            content_text = await extract_text_from_image(bytes(file_data), media_type)
        elif media_type == "excel":
            media_analysis = await analyze_excel_with_gpt(bytes(file_data), filename, context)
            content_text = await extract_csv_from_excel(bytes(file_data), filename)
        elif media_type == "word":
            media_analysis = await analyze_word_with_gpt(bytes(file_data), filename, context)
            content_text = await extract_text_from_word(bytes(file_data))
        elif media_type == "powerpoint":
            media_analysis = await analyze_pptx_with_gpt(bytes(file_data), filename, context)
            content_text = await extract_text_from_pptx(bytes(file_data))
        elif media_type == "video":
            media_analysis = await analyze_video_with_gemini(bytes(file_data), filename, context)
            # Для видео извлекаем транскрипт аудио
            import tempfile
            with tempfile.NamedTemporaryFile(delete=False, suffix='.mp4') as tmp:
                tmp.write(bytes(file_data))
                video_path = tmp.name
            try:
                content_text = await extract_transcript_from_audio(video_path)
            finally:
                if os.path.exists(video_path):
                    os.unlink(video_path)

    except Exception as e:
        logger.error(f"Ошибка обработки медиа: {e}")

    # Автоизвлечение фактов из анализа
    if media_analysis and len(media_analysis) > 50:
        try:
            await extract_and_save_facts(media_analysis, source=f"document:{filename or media_type_str}")
        except Exception as e:
            logger.debug(f"Fact extraction error: {e}")

    return media_type_str, media_analysis, content_text


def determine_message_type(message) -> tuple[str, str | None]:
    """Определяет тип сообщения и file_id если есть медиа."""
    if message.photo:
        return "photo", message.photo[-1].file_id
    elif message.video:
        return "video", message.video.file_id
    elif message.audio:
        return "audio", message.audio.file_id
    elif message.voice:
        return "voice", message.voice.file_id
    elif message.video_note:
        return "video_note", message.video_note.file_id
    elif message.document:
        return "document", message.document.file_id
    elif message.sticker:
        return "sticker", message.sticker.file_id
    elif message.animation:
        return "animation", message.animation.file_id
    elif message.location:
        return "location", None
    elif message.contact:
        return "contact", None
    elif message.poll:
        return "poll", None
    else:
        return "text", None


async def analyze_daily_documents(bot, chat_id: int, chat_title: str):
    """Анализирует все документы за день для указанного чата.

    Вызывается планировщиком в конце дня (23:55).
    """
    logger.info(f"Запуск анализа документов за день для чата {chat_title} ({chat_id})")

    table_name = ensure_table_exists(chat_id, chat_title)

    # Получаем все документы за сегодня без анализа
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            # Ищем документы за сегодня, которые не имеют анализа
            cur.execute(f"""
                SELECT message_id, media_file_id, message_text, message_type, timestamp
                FROM {table_name}
                WHERE DATE(timestamp) = CURRENT_DATE
                AND message_type IN ('pdf', 'excel', 'word', 'powerpoint', 'document', 'photo')
                AND (media_analysis IS NULL OR media_analysis = '')
                ORDER BY timestamp
            """)
            documents = cur.fetchall()
    finally:
        conn.close()

    if not documents:
        logger.info(f"Нет документов для анализа в чате {chat_title}")
        return

    logger.info(f"Найдено {len(documents)} документов для анализа в чате {chat_title}")

    # Формируем сводный анализ всех документов за день
    all_docs_info = []
    for idx, (msg_id, file_id, caption, msg_type, timestamp) in enumerate(documents, 1):
        time_str = timestamp.strftime("%H:%M")
        doc_info = f"{idx}. Документ {msg_type} в {time_str}"
        if caption:
            doc_info += f": {caption[:100]}"
        all_docs_info.append(doc_info)

    # Получаем контекст чата за день
    context = get_full_chat_context(table_name, chat_id, chat_title, 24)  # 24 часа

    # Создаем сводный анализ
    company_profile = get_company_profile()
    summary_prompt = f"""{company_profile}

Проанализируй все документы, отправленные в чат "{chat_title}" компании Фрумелад за сегодня.
Используй знания о компании из профиля выше.

=== СПИСОК ДОКУМЕНТОВ ЗА ДЕНЬ ===
{chr(10).join(all_docs_info)}

=== КОНТЕКСТ ЧАТА ЗА ДЕНЬ ===
{context}

Создай краткий сводный анализ всех документов:
1. Общая тематика документов
2. Ключевые данные и цифры
3. Важные моменты и выводы
4. Связь между документами (если есть)

Анализ должен быть структурированным и информативным."""

    try:
        if gpt_client:
            response = gpt_client.chat.completions.create(
                model="openai/gpt-4.1",
                max_tokens=6000,
                messages=[{"role": "user", "content": summary_prompt}],
            )
            summary_analysis = response.choices[0].message.content
            
            # Автоизвлечение фактов из сводного анализа дня
            try:
                await extract_and_save_facts(summary_analysis, source=f"daily_analysis:{chat_title}")
            except Exception as e:
                logger.debug(f"Fact extraction from daily analysis error: {e}")
            
            # Сохраняем сводный анализ в БД для последнего документа дня
            # (или можно создать отдельную таблицу для дневных отчетов)
            conn = get_db_connection()
            try:
                with conn.cursor() as cur:
                    # Обновляем последний документ дня со сводным анализом
                    last_msg_id = documents[-1][0]
                    cur.execute(f"""
                        UPDATE {table_name}
                        SET media_analysis = %s
                        WHERE message_id = %s
                    """, (f"📊 СВОДНЫЙ АНАЛИЗ ДОКУМЕНТОВ ЗА ДЕНЬ\n\n{summary_analysis}", last_msg_id))
                    conn.commit()
            finally:
                conn.close()

            # Отправляем сводный анализ в чат
            await bot.send_message(
                chat_id=chat_id,
                text=f"📊 *Сводный анализ документов за {datetime.now().strftime('%d.%m.%Y')}*\n\n"
                     f"Проанализировано документов: {len(documents)}\n\n"
                     f"{summary_analysis}",
                parse_mode="Markdown"
            )

            logger.info(f"Сводный анализ отправлен в чат {chat_title}")

    except Exception as e:
        logger.error(f"Ошибка при создании сводного анализа: {e}")


# ============================================================
# ОБРАБОТЧИКИ СООБЩЕНИЙ
# ============================================================

async def log_message(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик для логирования всех сообщений."""
    message = update.message or update.edited_message
    
    if not message or not message.chat:
        return
    
    if message.chat.type not in ["group", "supergroup"]:
        return
    
    chat_id = message.chat.id
    chat_title = message.chat.title or f"Chat_{abs(chat_id)}"
    
    table_name = ensure_table_exists(chat_id, chat_title)
    
    message_type, media_file_id = determine_message_type(message)

    # Проверяем, нужно ли отложить анализ для этой группы
    is_delayed_chat = chat_title == DELAYED_ANALYSIS_CHAT

    # Анализируем медиа если есть (кроме группы с отложенным анализом)
    media_analysis = ""
    content_text = ""
    if message.photo or message.video or message.voice or message.audio or (message.document and (message.document.mime_type or message.document.file_name)):
        # Для группы "Торты Отгрузки" не анализируем сразу - анализ будет в конце дня
        if not is_delayed_chat:
            analyzed_type, media_analysis, content_text = await download_and_analyze_media(context.bot, message, table_name)
            if analyzed_type != "media":
                message_type = analyzed_type
        else:
            # Для отложенного анализа просто определяем тип документа
            if message.document:
                doc = message.document
                filename = doc.file_name or ""
                filename_lower = filename.lower()
                if doc.mime_type == "application/pdf" or filename_lower.endswith(".pdf"):
                    message_type = "pdf"
                elif filename_lower.endswith(('.xlsx', '.xls')):
                    message_type = "excel"
                elif filename_lower.endswith(('.docx', '.doc')):
                    message_type = "word"
                elif filename_lower.endswith(('.pptx', '.ppt')):
                    message_type = "powerpoint"

            logger.info(f"Документ {message_type} отложен для анализа в конце дня (чат: {chat_title})")

        # Отправляем результат анализа в чат (только если анализ был выполнен)
        if media_analysis:
            try:
                # Формируем краткий анализ (2-3 строки)
                lines = [l for l in media_analysis.split('\n') if l.strip()]
                summary = '\n'.join(lines[:3])
                if len(summary) > 350:
                    summary = summary[:350] + "..."
                
                # В чат — только краткий анализ
                filename = ""
                if message.document and message.document.file_name:
                    filename = f" ({message.document.file_name})"
                
                # Создаём кнопку для получения полного анализа
                keyboard = InlineKeyboardMarkup([
                    [InlineKeyboardButton("📋 Полный анализ", callback_data=f"full_{message.message_id}")]
                ])
                
                await message.reply_text(f"📄 Анализ{filename}:\n\n{summary}", reply_markup=keyboard)
                
                # Рассылка полного анализа в личку тем, кто включил
                if len(media_analysis) > 400:  # Только если есть что добавить
                    chat_title = message.chat.title or "Чат"
                    sender_name = message.from_user.first_name if message.from_user else "Неизвестный"
                    
                    full_message = (
                        f"📄 *Полный анализ документа\n\n"
                        f"📍 Чат: {chat_title}\n"
                        f"👤 Отправил: {sender_name}\n"
                        f"📎 Файл: {filename.strip(' ()') or message_type}\n\n"
                        f"{media_analysis.replace('*', '✱').replace('_', '‗')}"
                    )
                    
                    # Получаем список всех пользователей с включённой рассылкой
                    conn = get_db_connection()
                    try:
                        with conn.cursor() as cur:
                            # Все пользователи с включённой рассылкой (независимо от чата)
                            cur.execute("""
                                SELECT user_id
                                FROM tg_full_analysis_settings
                                WHERE send_full_analysis = TRUE
                            """)
                            all_subscribers = [row[0] for row in cur.fetchall()]
                    finally:
                        conn.close()
                    
                    # Фильтруем: admin получает всё, остальные — только из своих чатов
                    users_to_notify = []
                    chat_id = message.chat.id
                    for uid in all_subscribers:
                        if uid == ADMIN_USER_ID:
                            users_to_notify.append(uid)
                        else:
                            try:
                                member = await context.bot.get_chat_member(chat_id=chat_id, user_id=uid)
                                if member.status in ('member', 'administrator', 'creator', 'restricted'):
                                    users_to_notify.append(uid)
                            except Exception:
                                pass  # не состоит в чате или ошибка
                    
                    # Отправляем в личку
                    for uid in users_to_notify:
                        try:
                            # Разбиваем если слишком длинный
                            if len(full_message) > 4000:
                                parts = [full_message[i:i+4000] for i in range(0, len(full_message), 4000)]
                                for i, part in enumerate(parts):
                                    await context.bot.send_message(
                                        chat_id=uid,
                                        text=part if i == 0 else f"...продолжение:\n\n{part}",
                                        parse_mode=None
                                    )
                            else:
                                await context.bot.send_message(
                                    chat_id=uid,
                                    text=full_message,
                                    parse_mode=None
                                )
                        except Exception as e:
                            logger.warning(f"Не удалось отправить анализ пользователю {uid}: {e}")
                            # Если бот заблокирован — отключаем рассылку
                            if "bot was blocked" in str(e).lower() or "chat not found" in str(e).lower():
                                set_user_analysis_setting(uid, "", "", False)
                    
            except Exception as e:
                logger.error(f"Ошибка отправки анализа: {e}")
    
    text = message.text or message.caption or ""
    
    message_data = {
        "message_id": message.message_id,
        "user_id": message.from_user.id if message.from_user else None,
        "username": message.from_user.username if message.from_user else None,
        "first_name": message.from_user.first_name if message.from_user else None,
        "last_name": message.from_user.last_name if message.from_user else None,
        "message_text": text,
        "message_type": message_type,
        "reply_to_message_id": message.reply_to_message.message_id if message.reply_to_message else None,
        "forward_from_user_id": None,
        "media_file_id": media_file_id,
        "media_analysis": media_analysis,
        "content_text": content_text,
        "timestamp": message.date
    }

    save_message(table_name, message_data)

    # Индексируем для векторного поиска (включая извлеченное содержимое)
    content_for_index = (message_data.get("message_text") or "") + " " + (message_data.get("media_analysis") or "") + " " + (message_data.get("content_text") or "")
    if content_for_index.strip():
        await index_new_message(table_name, message_data["message_id"], content_for_index.strip())
    logger.info(f"Сохранено сообщение {message.message_id} ({message_type}) в {table_name}")
    
    # Проверяем, есть ли у пользователя роль
    if message.from_user and ADMIN_USER_ID and message.from_user.id != ADMIN_USER_ID:
        user_role = get_user_role(message.from_user.id, chat_id)
        if not user_role:
            # Запрашиваем роль у администратора (только для новых пользователей)
            await request_user_role(context.bot, message, chat_id, chat_title)


async def request_user_role(bot, message, chat_id: int, chat_title: str):
    """Запрашивает роль пользователя у администратора."""
    if not ADMIN_USER_ID:
        return
    
    user = message.from_user
    user_name = f"{user.first_name or ''} {user.last_name or ''}".strip() or user.username or f"User_{user.id}"
    
    try:
        await bot.send_message(
            chat_id=ADMIN_USER_ID,
            text=f"👤 Новый пользователь без роли в чате \"{chat_title}\":\n\n"
                 f"Имя: {user_name}\n"
                 f"Username: @{user.username or 'нет'}\n"
                 f"ID: {user.id}\n\n"
                 f"Ответьте на это сообщение, указав роль пользователя.\n"
                 f"Например: Бухгалтер, Менеджер, Директор и т.д.",
            parse_mode="HTML"
        )
        
        # Сохраняем ожидание ответа
        pending_role_assignments[f"admin_{user.id}_{chat_id}"] = {
            "user_id": user.id,
            "chat_id": chat_id,
            "username": user.username,
            "first_name": user.first_name,
            "last_name": user.last_name
        }
    except Exception as e:
        logger.error(f"Ошибка запроса роли: {e}")


async def handle_admin_role_reply(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обрабатывает ответ администратора с ролью пользователя."""
    message = update.message
    
    if not message.reply_to_message or message.chat.type != "private":
        return
    
    if message.from_user.id != ADMIN_USER_ID:
        return
    
    # Ищем ожидающее назначение
    reply_text = message.reply_to_message.text or ""
    
    for key, pending in list(pending_role_assignments.items()):
        if key.startswith("admin_") and f"ID: {pending['user_id']}" in reply_text:
            role = message.text.strip()
            
            set_user_role(
                pending["user_id"],
                pending["chat_id"],
                role,
                pending["username"],
                pending["first_name"],
                pending["last_name"]
            )
            
            await message.reply_text(
                f"✅ Роль \"{role}\" назначена пользователю {pending['first_name'] or pending['username']}"
            )
            
            del pending_role_assignments[key]
            return


# ============================================================
# КОМАНДЫ
# ============================================================

async def start_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик команды /start. Поддерживает deep link /start element."""
    if update.message.chat.type == "private":
        # Deep link: /start element → сразу выдать данные Element X
        payload = context.args[0] if context.args else ""
        if payload == "element":
            # Делегируем в element_command
            await element_command(update, context)
            return

        await update.message.reply_text(
            "👋 Привет! Я корпоративный бот компании Фрумелад.\n\n"
            "📝 Сохраняю сообщения в базу знаний\n"
            "🖼 Анализирую документы через AI\n"
            "🔍 Отвечаю на вопросы по базе знаний\n\n"
            "Команды:\n"
            "/element - данные для входа в Element X\n"
            "/rooms - переотправить приглашения в комнаты Element\n"
            "/search <запрос> - поиск по базе знаний\n"
            "/roles - показать пользователей без ролей\n"
            "/stats - статистика чата\n"
            "/analysis - настройка рассылки анализа документов"
        )
    else:
        await update.message.reply_text(
            "✅ Бот активирован в этом чате.\n"
            "📝 Логирование сообщений\n"
            "🖼 Анализ документов с контекстом\n"
            "👥 Учёт ролей участников\n\n"
            "Команды: /roles, /stats, /search, /analysis"
        )


async def roles_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает пользователей без ролей."""
    message = update.message
    
    if message.chat.type not in ["group", "supergroup"]:
        await message.reply_text("Эта команда работает только в групповых чатах.")
        return
    
    chat_id = message.chat.id
    chat_title = message.chat.title or f"Chat_{abs(chat_id)}"
    table_name = sanitize_table_name(chat_id, chat_title)
    
    users_without_roles = get_users_without_roles(chat_id, table_name)
    
    if not users_without_roles:
        # Показываем всех пользователей с ролями
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    SELECT first_name, last_name, role 
                    FROM tg_user_roles 
                    WHERE chat_id = %s 
                    ORDER BY role, first_name
                """, (chat_id,))
                users_with_roles = cur.fetchall()
        finally:
            conn.close()
        
        if users_with_roles:
            response = "✅ Все пользователи имеют роли:\n\n"
            for first_name, last_name, role in users_with_roles:
                name = f"{first_name or ''} {last_name or ''}".strip() or "Без имени"
                response += f"• {name} — {role}\n"
            await message.reply_text(response)
        else:
            await message.reply_text("В этом чате пока нет пользователей с назначенными ролями.")
        return
    
    response = "👥 **Пользователи без ролей:**\n\n"
    for i, (user_id, username, first_name, last_name) in enumerate(users_without_roles, 1):
        name = f"{first_name or ''} {last_name or ''}".strip() or username or f"User_{user_id}"
        response += f"{i}. {name} (@{username or 'нет'})\n"
    
    response += "\n**Чтобы назначить роль**, ответьте на это сообщение в формате:\n"
    response += "`1 Директор`\n`2 Бухгалтер`\n`3 Менеджер`"
    
    sent_message = await message.reply_text(response, parse_mode="Markdown")
    
    pending_role_assignments[chat_id] = {
        "message_id": sent_message.message_id,
        "users": users_without_roles,
        "table_name": table_name
    }


async def handle_role_assignment(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обрабатывает назначение ролей в групповом чате."""
    message = update.message
    if not message:
        return

    if not message.reply_to_message:
        return
    
    chat_id = message.chat.id
    
    # Обработка в личных сообщениях (ответ администратора)
    if message.chat.type == "private":
        await handle_admin_role_reply(update, context)
        return
    
    # Обработка в групповом чате
    if chat_id not in pending_role_assignments:
        return
    
    pending = pending_role_assignments[chat_id]
    
    if message.reply_to_message.message_id != pending["message_id"]:
        return
    
    lines = message.text.strip().split('\n')
    assigned = []
    
    for line in lines:
        parts = line.strip().split(' ', 1)
        if len(parts) != 2:
            continue
        
        try:
            index = int(parts[0]) - 1
            role = parts[1].strip()
            
            if 0 <= index < len(pending["users"]):
                user_id, username, first_name, last_name = pending["users"][index]
                set_user_role(user_id, chat_id, role, username, first_name, last_name)
                name = f"{first_name or ''} {last_name or ''}".strip() or username or f"User_{user_id}"
                assigned.append(f"{name} → {role}")
        except (ValueError, IndexError):
            continue
    
    if assigned:
        response = "✅ **Роли назначены:**\n" + "\n".join(assigned)
        await message.reply_text(response, parse_mode="Markdown")
        
        remaining = get_users_without_roles(chat_id, pending["table_name"])
        if remaining:
            await message.reply_text(f"Осталось пользователей без ролей: {len(remaining)}\nИспользуйте /roles чтобы продолжить.")
        else:
            del pending_role_assignments[chat_id]
    else:
        await message.reply_text("Не удалось распознать. Используйте формат:\n`1 Директор`", parse_mode="Markdown")


async def stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает статистику по текущему чату."""
    message = update.message
    
    if message.chat.type not in ["group", "supergroup"]:
        await message.reply_text("Эта команда работает только в групповых чатах.")
        return
    
    chat_id = message.chat.id
    chat_title = message.chat.title or f"Chat_{abs(chat_id)}"
    table_name = sanitize_table_name(chat_id, chat_title)
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT EXISTS (
                    SELECT FROM information_schema.tables 
                    WHERE table_name = %s
                )
            """, (table_name,))
            
            if not cur.fetchone()[0]:
                await message.reply_text("📊 Пока нет данных для этого чата.")
                return
            
            cur.execute(sql.SQL("SELECT COUNT(*) FROM {}").format(sql.Identifier(table_name)))
            total_messages = cur.fetchone()[0]
            
            cur.execute(sql.SQL("SELECT COUNT(DISTINCT user_id) FROM {}").format(sql.Identifier(table_name)))
            unique_users = cur.fetchone()[0]
            
            cur.execute(sql.SQL("""
                SELECT COUNT(*) FROM {} WHERE message_type IN ('photo', 'pdf', 'image', 'excel', 'word', 'powerpoint', 'video')
            """).format(sql.Identifier(table_name)))
            media_count = cur.fetchone()[0]
            
            cur.execute(sql.SQL("""
                SELECT COUNT(*) FROM {} WHERE media_analysis IS NOT NULL AND media_analysis != ''
            """).format(sql.Identifier(table_name)))
            analyzed_count = cur.fetchone()[0]
            
            cur.execute("""
                SELECT COUNT(*) FROM tg_user_roles WHERE chat_id = %s AND role IS NOT NULL
            """, (chat_id,))
            roles_count = cur.fetchone()[0]
            
            cur.execute(sql.SQL("""
                SELECT COUNT(*) FROM {} WHERE timestamp::date = CURRENT_DATE
            """).format(sql.Identifier(table_name)))
            today_messages = cur.fetchone()[0]
            
            stats_text = (
                f"📊 **Статистика чата**\n\n"
                f"📝 Всего сообщений: {total_messages:,}\n"
                f"👥 Участников: {unique_users}\n"
                f"🏷 С ролями: {roles_count}\n"
                f"📅 Сегодня: {today_messages}\n"
                f"📎 Медиафайлов: {media_count}\n"
                f"🤖 Проанализировано: {analyzed_count}"
            )
            
            await message.reply_text(stats_text, parse_mode="Markdown")
            
    finally:
        conn.close()


async def search_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Поиск по сообщениям в текущем чате."""
    message = update.message
    
    if message.chat.type not in ["group", "supergroup"]:
        await message.reply_text("Эта команда работает только в групповых чатах.")
        return
    
    query = ' '.join(context.args) if context.args else None
    
    if not query:
        await message.reply_text(
            "🔍 Использование: /search <запрос>\n\n"
            "Пример: /search накладная сахар"
        )
        return
    
    chat_id = message.chat.id
    chat_title = message.chat.title or f"Chat_{abs(chat_id)}"
    table_name = sanitize_table_name(chat_id, chat_title)
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(sql.SQL("""
                SELECT 
                    m.timestamp,
                    COALESCE(m.first_name, m.username, 'Неизвестный') as author,
                    r.role,
                    m.message_type,
                    LEFT(COALESCE(m.message_text, '') || ' ' || COALESCE(m.media_analysis, ''), 300) as content
                FROM {} m
                LEFT JOIN tg_user_roles r ON m.user_id = r.user_id AND r.chat_id = %s
                WHERE to_tsvector('russian', COALESCE(m.message_text, '') || ' ' || COALESCE(m.media_analysis, '')) 
                      @@ plainto_tsquery('russian', %s)
                ORDER BY m.timestamp DESC
                LIMIT 10
            """).format(sql.Identifier(table_name)), (chat_id, query,))
            
            results = cur.fetchall()
            
            if not results:
                await message.reply_text(f"🔍 По запросу «{query}» ничего не найдено.")
                return
            
            response = f"🔍 **Результаты поиска:** «{query}»\n\n"
            for ts, author, role, msg_type, content in results:
                date_str = ts.strftime("%d.%m.%Y %H:%M")
                role_str = f" [{role}]" if role else ""
                type_emoji = {"photo": "🖼", "pdf": "📄", "document": "📎", "excel": "📊", "word": "📝", "powerpoint": "📽", "video": "🎬"}.get(msg_type, "💬")
                content_preview = content[:150] + "..." if len(content) > 150 else content
                response += f"{type_emoji} {date_str} | **{author}**{role_str}\n{content_preview}\n\n"
            
            await message.reply_text(response, parse_mode="Markdown")
            
    finally:
        conn.close()


async def analysis_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /analysis — управление рассылкой полного анализа документов в личку."""
    user = update.effective_user
    user_id = user.id
    username = user.username or ""
    first_name = user.first_name or ""
    
    args = context.args
    
    if not args:
        current = get_user_analysis_setting(user_id)
        status = "✅ включена" if current else "❌ выключена"
        await update.message.reply_text(
            f"📄 *Рассылка полного анализа документов:* {status}\n\n"
            f"Когда кто-то отправляет документ в чат, бот анализирует его.\n"
            f"В чат приходит краткий анализ (2-3 строки).\n"
            f"Если включено — полный анализ приходит вам в личку.\n\n"
            f"Команды:\n"
            f"`/analysis on` — включить\n"
            f"`/analysis off` — выключить",
            parse_mode="Markdown"
        )
        return
    
    action = args[0].lower()
    
    if action == 'on':
        set_user_analysis_setting(user_id, username, first_name, True)
        await update.message.reply_text(
            "✅ Готово! Теперь полный анализ документов будет приходить вам в личные сообщения.\n\n"
            "⚠️ Убедитесь, что вы начали диалог с ботом (напишите /start в личку боту)."
        )
        logger.info(f"Пользователь {first_name} ({user_id}) включил рассылку анализа")
    
    elif action == 'off':
        set_user_analysis_setting(user_id, username, first_name, False)
        await update.message.reply_text("❌ Рассылка полного анализа отключена.")
        logger.info(f"Пользователь {first_name} ({user_id}) отключил рассылку анализа")
    
    else:
        await update.message.reply_text(
            "Используйте:\n"
            "`/analysis on` — включить\n"
            "`/analysis off` — выключить",
            parse_mode="Markdown"
        )


async def analysis_list_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /analysis_list — показать кто включил рассылку (только для админа)."""
    user_id = update.effective_user.id
    
    if user_id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Эта команда доступна только администратору.")
        return
    
    users = get_users_with_full_analysis_enabled()
    
    if not users:
        await update.message.reply_text("📭 Никто не включил рассылку полного анализа.")
        return
    
    response = "📄 *Пользователи с включённой рассылкой анализа:*\n\n"
    for uid, username, first_name in users:
        name = first_name or username or str(uid)
        username_str = f" (@{username})" if username else ""
        response += f"• {name}{username_str}\n"
    
    response += f"\n*Всего:* {len(users)}"
    
    await update.message.reply_text(response, parse_mode="Markdown")
    

async def chats_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает список всех логируемых чатов."""
    if update.message.chat.type != "private":
        await update.message.reply_text("Эта команда доступна только в личных сообщениях.")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT chat_title, total_messages, last_message_at
                FROM tg_chats_metadata
                ORDER BY last_message_at DESC NULLS LAST
                LIMIT 20
            """)
            chats = cur.fetchall()
            
            if not chats:
                await update.message.reply_text("📭 Пока нет подключенных чатов.")
                return
            
            response = "📋 **Логируемые чаты:**\n\n"
            for title, total, last_msg in chats:
                last_str = last_msg.strftime("%d.%m.%Y %H:%M") if last_msg else "—"
                total = total or 0
                response += f"• **{title}**\n  Сообщений: {total:,} | {last_str}\n\n"
            
            await update.message.reply_text(response, parse_mode="Markdown")
            
    finally:
        conn.close()


# ============================================================
# ЗАПУСК
# ============================================================

async def error_handler(update: object, context: ContextTypes.DEFAULT_TYPE) -> None:
    """Обработчик ошибок — логирует и продолжает работу."""
    import traceback
    
    logger.error(f"Ошибка при обработке обновления: {context.error}")
    
    # Логируем traceback
    tb_list = traceback.format_exception(None, context.error, context.error.__traceback__)
    tb_string = "".join(tb_list)
    logger.error(f"Traceback:\n{tb_string[:1000]}")
    
    # Отправляем алерт администратору (не чаще раза в час)
    if ADMIN_USER_ID:
        try:
            error_text = str(context.error)[:200]
            await context.bot.send_message(
                chat_id=ADMIN_USER_ID,
                text=f"⚠️ Ошибка бота:\n\n{error_text}"
            )
        except:
            pass  # Не падаем если не можем отправить алерт

# ============================================================
# RAG АГЕНТ В ЛИЧНЫХ СООБЩЕНИЯХ
# ============================================================

async def handle_private_rag(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик личных сообщений — RAG агент без @упоминания.
    
    В личном чате с ботом можно просто писать вопросы без @упоминания.
    """
    message = update.message
    if not message or not message.text:
        return
    
    # Работаем только в личном чате
    if message.chat.type != "private":
        return
    
    question = message.text.strip()
    
    if not question:
        return
    
    # Пропускаем простые приветствия — на них ответит /start
    if question.lower() in ['привет', 'hi', 'hello', 'старт', 'start']:
        await message.reply_text(
            "👋 Привет! Я RAG-агент компании.\n\n"
            "Задайте вопрос — я поищу ответ в базе данных:\n"
            "• 1С (цены, закупки, номенклатура)\n"
            "• Telegram чаты\n"
            "• Email переписка\n\n"
            "Примеры:\n"
            "• по какой цене покупали сахар за последний месяц?\n"
            "• что обсуждали про поставщика Агросервер?\n"
            "• какие были проблемы с доставкой на прошлой неделе?"
        )
        return
    
    # Отправляем индикатор "печатает"
    await context.bot.send_chat_action(chat_id=message.chat.id, action="typing")
    
    try:
        # Reply-chain: если пользователь реплайнул на ответ бота — подхватываем
        # предыдущий Q/A из chat_data (PTB).
        prev_context = None
        if message.reply_to_message and message.reply_to_message.from_user \
                and message.reply_to_message.from_user.is_bot:
            replied_text = (message.reply_to_message.text or "")[:600]
            history = context.chat_data.get("rag_history", [])
            for h in reversed(history):
                if h.get("answer", "")[:300] == replied_text[:300] \
                        or replied_text.startswith(h.get("answer", "")[:200]):
                    prev_context = {"question": h["question"], "answer": h["answer"]}
                    logger.info(f"Reply-chain matched prev Q: '{prev_context['question'][:60]}'")
                    break

        response = await process_rag_query(question, "", user_info={
            "user_id": message.from_user.id,
            "username": message.from_user.username,
            "first_name": message.from_user.first_name,
            "chat_id": message.chat.id,
            "chat_type": "private",
        }, prev_context=prev_context)

        # Сохраняем Q/A в историю для будущих reply
        hist = context.chat_data.setdefault("rag_history", [])
        hist.append({"question": question, "answer": response})
        context.chat_data["rag_history"] = hist[-10:]  # последние 10

        # Отправляем ответ
        if len(response) > 4000:
            parts = [response[i:i+4000] for i in range(0, len(response), 4000)]
            for part in parts:
                await message.reply_text(part)
        else:
            await message.reply_text(response)

        logger.info(f"RAG ответ в личку: {len(response)} символов (follow-up={prev_context is not None})")

    except Exception as e:
        logger.error(f"Ошибка RAG агента в личном чате: {e}")
        await message.reply_text("Произошла ошибка при обработке запроса. Попробуйте позже.")
        

async def handle_mention(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик упоминания бота — RAG агент."""
    message = update.message
    if not message or not message.text:
        return
    
    # Проверяем, упомянут ли бот
    bot_username = (await context.bot.get_me()).username
    
    if f"@{bot_username}" not in message.text:
        return
    
    # Извлекаем вопрос (убираем упоминание бота)
    question = message.text.replace(f"@{bot_username}", "").strip()
    
    if not question:
        await message.reply_text("Задайте вопрос после упоминания бота.\n\nПример: @имя_бота какой курс доллара?")
        return
    
    # Отправляем индикатор "печатает"
    await context.bot.send_chat_action(chat_id=message.chat.id, action="typing")
    
    try:
        # Reply-chain (как в личке)
        prev_context = None
        if message.reply_to_message and message.reply_to_message.from_user \
                and message.reply_to_message.from_user.is_bot:
            replied_text = (message.reply_to_message.text or "")[:600]
            history = context.chat_data.get("rag_history", [])
            for h in reversed(history):
                if h.get("answer", "")[:300] == replied_text[:300] \
                        or replied_text.startswith(h.get("answer", "")[:200]):
                    prev_context = {"question": h["question"], "answer": h["answer"]}
                    logger.info(f"Reply-chain (mention) matched prev Q: '{prev_context['question'][:60]}'")
                    break

        response = await process_rag_query(question, "", user_info={
            "user_id": message.from_user.id,
            "username": message.from_user.username,
            "first_name": message.from_user.first_name,
            "chat_id": message.chat.id,
            "chat_type": message.chat.type,
        }, prev_context=prev_context)

        hist = context.chat_data.setdefault("rag_history", [])
        hist.append({"question": question, "answer": response})
        context.chat_data["rag_history"] = hist[-10:]

        if len(response) > 4000:
            parts = [response[i:i+4000] for i in range(0, len(response), 4000)]
            for part in parts:
                await message.reply_text(part)
        else:
            await message.reply_text(response)

        logger.info(f"RAG ответ отправлен: {len(response)} символов (follow-up={prev_context is not None})")
        
    except Exception as e:
        logger.error(f"Ошибка RAG агента: {e}")
        await message.reply_text(f"Произошла ошибка при обработке запроса. Попробуйте позже.")

# ============================================================
# EMAIL LOGGER: ВСПОМОГАТЕЛЬНЫЕ ФУНКЦИИ
# ============================================================

def format_email_age(dt) -> str:
    """Форматирует возраст для email."""
    if not dt:
        return "никогда"
    
    from datetime import datetime
    if dt.tzinfo:
        dt = dt.replace(tzinfo=None)
    
    delta = datetime.now() - dt
    
    if delta.days > 30:
        return f"{delta.days // 30} мес."
    elif delta.days > 0:
        return f"{delta.days} дн."
    elif delta.seconds > 3600:
        return f"{delta.seconds // 3600} ч."
    elif delta.seconds > 60:
        return f"{delta.seconds // 60} мин."
    else:
        return "сейчас"


def truncate_text(text: str, max_len: int = 100) -> str:
    """Обрезает текст."""
    if not text:
        return ""
    if len(text) <= max_len:
        return text
    return text[:max_len-3] + "..."


def format_thread_status(lifecycle_status: str, resolution_outcome: str | None) -> str:
    """Форматирует 2-слойный статус ветки для отображения в боте."""
    if lifecycle_status == "open":
        return "📬 Открыта"
    if lifecycle_status == "pending_resolution":
        return "⏳ Ожидает подтверждения"
    if lifecycle_status == "archived":
        return "📦 В архиве"
    if lifecycle_status == "closed":
        if resolution_outcome == "cancelled":
            return "❌ Закрыта (отменено)"
        if resolution_outcome == "resolved":
            return "✅ Закрыта (решено)"
        return "📧 Закрыта"
    return lifecycle_status or "unknown"


# ============================================================
# EMAIL LOGGER: КОМАНДЫ
# ============================================================

async def open_threads_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает открытые ветки email переписки."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Нет доступа к этой команде")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT 
                    t.id,
                    t.subject_normalized,
                    t.message_count,
                    t.last_message_at,
                    t.priority,
                    t.lifecycle_status,
                    t.resolution_outcome
                FROM email_threads t
                WHERE t.lifecycle_status IN ('open', 'pending_resolution')
                ORDER BY 
                    CASE t.priority 
                        WHEN 'high' THEN 1 
                        WHEN 'medium' THEN 2 
                        ELSE 3 
                    END,
                    t.last_message_at DESC
                LIMIT 20
            """)
            threads = cur.fetchall()
    except Exception as e:
        logger.error(f"Ошибка получения веток: {e}")
        await update.message.reply_text(
            "❌ Таблицы email логгера не найдены.\n\n"
            "Примените миграции:\n"
            "`psql -d knowledge_base -f 001_init_email_logger.sql`\n"
            "`psql -d knowledge_base -f 004_email_thread_status_model.sql`",
            parse_mode="Markdown"
        )
        return
    finally:
        conn.close()
    
    if not threads:
        await update.message.reply_text("✅ Нет открытых веток email переписки")
        return
    
    text = "📬 *Открытые ветки переписки:*\n\n"
    
    for thread_id, subject, msg_count, last_msg_at, priority, lifecycle_status, resolution_outcome in threads:
        priority_icon = {'high': '🔴', 'medium': '🟡', 'low': '🟢'}.get(priority or 'medium', '⚪')
        status_icon = '⏳' if lifecycle_status == 'pending_resolution' else '📨'
        age = format_email_age(last_msg_at)
        subject_short = truncate_text(subject or "Без темы", 45)
        
        text += (
            f"{priority_icon}{status_icon} *{subject_short}*\n"
            f"   📨 {msg_count or 0} писем • {age}\n"
            f"   /emailthread\\_{thread_id}\n\n"
        )
    
    await update.message.reply_text(text, parse_mode="Markdown")


async def show_email_thread_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает детали ветки по команде /emailthread_N."""
    import json
    
    text = update.message.text
    match = re.search(r'/emailthread_(\d+)', text)
    if not match:
        await update.message.reply_text("❌ Укажите ID ветки: /emailthread_123")
        return
    
    thread_id = int(match.group(1))
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT 
                    id, subject_normalized, message_count, last_message_at,
                    priority, lifecycle_status, resolution_outcome, participant_emails, topic_tags,
                    summary_short, key_decisions, action_items
                FROM email_threads WHERE id = %s
            """, (thread_id,))
            row = cur.fetchone()
            
            if not row:
                await update.message.reply_text("❌ Ветка не найдена")
                return
            
            (tid, subject, msg_count, last_msg_at, priority, lifecycle_status, resolution_outcome,
             participants, tags, summary, decisions, actions) = row
             
            # Получаем последние сообщения
            cur.execute("""
                SELECT from_address, body_text, received_at
                FROM email_messages
                WHERE thread_id = %s
                ORDER BY received_at DESC
                LIMIT 3
            """, (thread_id,))
            messages = cur.fetchall()
    finally:
        conn.close()
    
    # Статус
    status_str = format_thread_status(lifecycle_status, resolution_outcome)
    
    # Приоритет
    priority_map = {'high': '🔴 Высокий', 'medium': '🟡 Средний', 'low': '🟢 Низкий'}
    priority_str = priority_map.get(priority, priority or 'medium')
    
    # Формируем ответ
    response = (
        f"📧 *{truncate_text(subject or 'Без темы', 50)}*\n\n"
        f"*Статус:* {status_str}\n"
        f"*Приоритет:* {priority_str}\n"
        f"*Сообщений:* {msg_count or 0}\n"
        f"*Последнее:* {format_email_age(last_msg_at)}\n"
    )
    
    if participants:
        p_list = participants[:3] if isinstance(participants, list) else []
        if p_list:
            response += f"*Участники:* {', '.join(p_list)}\n"
    
    if tags and isinstance(tags, list):
        response += f"*Теги:* {', '.join(tags)}\n"
    
    if summary:
        response += f"\n📝 *Саммари:*\n{summary}\n"
        
        if decisions and isinstance(decisions, list):
            response += "\n*Решения:*\n"
            for d in decisions[:5]:
                response += f"✓ {d}\n"
        
        if actions:
            items = actions if isinstance(actions, list) else json.loads(actions) if isinstance(actions, str) else []
            if items:
                response += "\n*Задачи:*\n"
                for item in items[:5]:
                    if isinstance(item, dict):
                        assignee = item.get('assignee', '?')
                        task = item.get('task', '')
                        response += f"• {assignee}: {task}\n"
    
    # Последние сообщения
    if messages:
        response += "\n📜 *Последние сообщения:*\n"
        for from_addr, body, received_at in messages:
            date_str = received_at.strftime('%d.%m %H:%M') if received_at else ""
            body_short = truncate_text(body or "", 150)
            response += f"\n_{from_addr}_ ({date_str}):\n{body_short}\n"
    
    # Кнопки
    from telegram import InlineKeyboardMarkup, InlineKeyboardButton
    keyboard = InlineKeyboardMarkup([
        [
            InlineKeyboardButton("✅ Решена", callback_data=f"email_resolve:{thread_id}"),
            InlineKeyboardButton("📦 Архив", callback_data=f"email_archive:{thread_id}"),
        ]
    ])
    
    await update.message.reply_text(response[:4000], parse_mode="Markdown", reply_markup=keyboard)


async def email_callback_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик callback-кнопок для email."""
    query = update.callback_query
    await query.answer()
    
    data = query.data
    
    if data.startswith("email_resolve:"):
        thread_id = int(data.split(":")[1])
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    UPDATE email_threads
                    SET lifecycle_status = 'closed',
                        resolution_outcome = 'resolved',
                        status = 'resolved',
                        resolution_confirmed = true,
                        updated_at = NOW()
                    WHERE id = %s
                """, (thread_id,))
                conn.commit()
        finally:
            conn.close()
        await query.answer("✅ Ветка отмечена как решённая")
        await query.edit_message_reply_markup(reply_markup=None)
    
    elif data.startswith("email_archive:"):
        thread_id = int(data.split(":")[1])
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    UPDATE email_threads
                    SET lifecycle_status = 'archived',
                        status = 'archived',
                        updated_at = NOW()
                    WHERE id = %s
                """, (thread_id,))
                conn.commit()
        finally:
            conn.close()
        await query.answer("📦 Ветка перемещена в архив")
        await query.edit_message_reply_markup(reply_markup=None)


async def email_stats_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает статистику email логгера."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Нет доступа")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            # Проверяем существование таблиц
            cur.execute("""
                SELECT EXISTS (
                    SELECT FROM information_schema.tables 
                    WHERE table_name = 'email_messages'
                )
            """)
            if not cur.fetchone()[0]:
                await update.message.reply_text(
                    "📊 Email логгер ещё не настроен.\n\n"
                    "Примените миграцию:\n"
                    "`psql -d knowledge_base -f 001_init_email_logger.sql`",
                    parse_mode="Markdown"
                )
                return
            
            cur.execute("""
                SELECT
                    (SELECT COUNT(*) FROM monitored_mailboxes WHERE is_active = true),
                    (SELECT COUNT(*) FROM email_messages),
                    (SELECT COUNT(*) FROM email_threads),
                    (SELECT COUNT(*) FROM email_threads WHERE lifecycle_status IN ('open', 'pending_resolution')),
                    (SELECT COUNT(*) FROM email_attachments),
                    (SELECT COUNT(*) FROM email_attachments WHERE analysis_status = 'pending')
            """)
            mailboxes, messages, threads, open_threads, attachments, pending = cur.fetchone()
            
            cur.execute("""
                SELECT email, last_sync_at, sync_status
                FROM monitored_mailboxes
                WHERE last_sync_at IS NOT NULL
                ORDER BY last_sync_at DESC
                LIMIT 1
            """)
            last_sync = cur.fetchone()
            
            cur.execute("""
                SELECT COUNT(*) FROM monitored_mailboxes WHERE sync_status = 'error'
            """)
            error_count = cur.fetchone()[0]
    finally:
        conn.close()
    
    text = (
        "📊 *Статистика Email Логгера:*\n\n"
        f"📬 Ящиков: {mailboxes or 0}\n"
        f"📨 Сообщений: {messages or 0:,}\n"
        f"🔗 Веток: {threads or 0} (открытых: {open_threads or 0})\n"
        f"📎 Вложений: {attachments or 0} (в очереди: {pending or 0})\n"
    )
    
    if last_sync:
        email, sync_at, status = last_sync
        text += f"\n*Последняя синхронизация:*\n{email} — {format_email_age(sync_at)}\n"
    
    if error_count:
        text += f"\n⚠️ Ящиков с ошибками: {error_count}"
    
    await update.message.reply_text(text, parse_mode="Markdown")


async def sync_status_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает статус синхронизации ящиков."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT email, sync_status, last_sync_at
                FROM monitored_mailboxes
                WHERE is_active = true
                ORDER BY last_sync_at DESC NULLS LAST
                LIMIT 30
            """)
            mailboxes = cur.fetchall()
    except Exception as e:
        await update.message.reply_text(f"❌ Ошибка: {e}")
        return
    finally:
        conn.close()
    
    if not mailboxes:
        await update.message.reply_text("📬 Нет активных почтовых ящиков")
        return
    
    status_icons = {'idle': '✅', 'syncing': '🔄', 'initial_load': '📥', 'error': '❌'}
    
    text = "📬 *Статус синхронизации:*\n\n"
    
    for email, status, last_sync in mailboxes:
        icon = status_icons.get(status or 'idle', '❓')
        age = format_email_age(last_sync) if last_sync else "—"
        mailbox_name = email.split('@')[0] if email else "?"
        text += f"{icon} `{mailbox_name}` {age}\n"
    
    await update.message.reply_text(text[:4000], parse_mode="Markdown")


async def search_email_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Поиск по email сообщениям."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Нет доступа")
        return
    
    if not context.args:
        await update.message.reply_text(
            "🔍 *Поиск по email:*\n\n"
            "`/search_email накладная сахар`",
            parse_mode="Markdown"
        )
        return
    
    query_text = ' '.join(context.args)
    
    if len(query_text) < 3:
        await update.message.reply_text("Запрос слишком короткий")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT 
                    m.subject,
                    m.from_address,
                    m.received_at,
                    t.id as thread_id
                FROM email_messages m
                LEFT JOIN email_threads t ON t.id = m.thread_id
                WHERE 
                    m.subject ILIKE %s OR
                    m.body_text ILIKE %s OR
                    m.from_address ILIKE %s
                ORDER BY m.received_at DESC
                LIMIT 10
            """, (f"%{query_text}%", f"%{query_text}%", f"%{query_text}%"))
            results = cur.fetchall()
    finally:
        conn.close()
    
    if not results:
        await update.message.reply_text(f"❌ По запросу «{query_text}» ничего не найдено")
        return
    
    text = f"🔍 *Результаты «{query_text}»:*\n\n"
    
    for subject, from_addr, received_at, thread_id in results:
        subject_short = truncate_text(subject or "Без темы", 40)
        date = received_at.strftime('%d.%m.%Y') if received_at else ""
        thread_link = f"/emailthread\\_{thread_id}" if thread_id else ""
        
        text += f"📧 *{subject_short}*\n"
        text += f"   {from_addr or '?'} • {date}\n"
        if thread_link:
            text += f"   {thread_link}\n"
        text += "\n"
    
    await update.message.reply_text(text[:4000], parse_mode="Markdown")


async def add_employee_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Добавляет сотрудника."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return
    
    if not context.args:
        await update.message.reply_text(
            "👤 *Добавление сотрудника:*\n\n"
            "`/add_employee Иванов Иван | Бухгалтерия | Бухгалтер`\n"
            "`/add_employee Петрова Мария | Производство`\n"
            "`/add_employee Сидоров Пётр`",
            parse_mode="Markdown"
        )
        return
    
    full_text = ' '.join(context.args)
    parts = [p.strip() for p in full_text.split('|')]
    
    full_name = parts[0] if len(parts) > 0 else None
    department = parts[1] if len(parts) > 1 else None
    position = parts[2] if len(parts) > 2 else None
    
    if not full_name:
        await update.message.reply_text("❌ Укажите имя")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                INSERT INTO employees (full_name, department, position, is_active)
                VALUES (%s, %s, %s, true)
                RETURNING id
            """, (full_name, department, position))
            emp_id = cur.fetchone()[0]
            conn.commit()
    except Exception as e:
        await update.message.reply_text(f"❌ Ошибка: {e}")
        return
    finally:
        conn.close()
    
    text = f"✅ *Сотрудник добавлен:*\n\n👤 {full_name}\n"
    if department:
        text += f"🏢 {department}\n"
    if position:
        text += f"💼 {position}\n"
    text += f"\nID: {emp_id}"
    
    await update.message.reply_text(text, parse_mode="Markdown")


async def assign_email_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Назначает email сотруднику."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return
    
    if len(context.args) < 2:
        await update.message.reply_text(
            "📧 *Назначение email:*\n\n"
            "`/assign_email <ID сотрудника> <email>`\n\n"
            "Пример:\n"
            "`/assign_email 1 accountant@totsamiy.com`\n\n"
            "Список сотрудников: /list\\_employees",
            parse_mode="Markdown"
        )
        return
    
    try:
        employee_id = int(context.args[0])
        email = context.args[1].lower()
    except:
        await update.message.reply_text("❌ Формат: /assign_email <ID> <email>")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            # Проверяем сотрудника
            cur.execute("SELECT full_name FROM employees WHERE id = %s", (employee_id,))
            emp = cur.fetchone()
            if not emp:
                await update.message.reply_text(f"❌ Сотрудник ID {employee_id} не найден")
                return
            
            # Назначаем email
            cur.execute("""
                INSERT INTO employee_emails (employee_id, email, is_primary, assigned_by)
                VALUES (%s, %s, true, %s)
                ON CONFLICT (employee_id, email) DO NOTHING
            """, (employee_id, email, update.effective_user.id))
            conn.commit()
    finally:
        conn.close()
    
    await update.message.reply_text(
        f"✅ *Email назначен:*\n\n"
        f"👤 {emp[0]}\n"
        f"📧 {email}",
        parse_mode="Markdown"
    )


async def list_employees_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает список сотрудников."""
    
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Нет доступа")
        return
    
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT e.id, e.full_name, e.department, 
                       array_agg(ee.email) FILTER (WHERE ee.email IS NOT NULL) as emails
                FROM employees e
                LEFT JOIN employee_emails ee ON ee.employee_id = e.id
                WHERE e.is_active = true
                GROUP BY e.id, e.full_name, e.department
                ORDER BY e.full_name
                LIMIT 30
            """)
            employees = cur.fetchall()
    except Exception as e:
        await update.message.reply_text(f"❌ Ошибка: {e}")
        return
    finally:
        conn.close()
    
    if not employees:
        await update.message.reply_text("👤 Нет сотрудников. Добавьте через /add\\_employee", parse_mode="Markdown")
        return
    
    text = "👥 *Сотрудники:*\n\n"
    
    for emp_id, name, dept, emails in employees:
        dept_str = f" ({dept})" if dept else ""
        email_str = f"\n   📧 {', '.join(emails)}" if emails and emails[0] else ""
        text += f"*{emp_id}.* {name}{dept_str}{email_str}\n"
    
    await update.message.reply_text(text[:4000], parse_mode="Markdown")

BOM_GROUP_CHAT_ID = -1003559489741  # Группа "Новые продукты и конкуренты"
BOM_SERVER_URL = "http://95.174.92.209"

async def bom_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /bom — генерация ссылки на отчёт BOM"""
    user = update.effective_user
    if not user:
        return

    # Проверяем: команда только в личном чате
    if update.effective_chat.type != 'private':
        await update.message.reply_text("Эта команда работает только в личном чате с ботом.")
        return

    # Проверяем членство в группе
    try:
        member = await context.bot.get_chat_member(chat_id=BOM_GROUP_CHAT_ID, user_id=user.id)
        if member.status in ('left', 'kicked'):
            await update.message.reply_text("⛔ У вас нет доступа к отчёту BOM.\nОбратитесь к администратору.")
            return
    except Exception as e:
        logger.warning(f"Ошибка проверки членства для /bom: {e}")
        if str(user.id) != ADMIN_USER_ID:
            await update.message.reply_text("⛔ Не удалось проверить доступ. Попробуйте позже.")
            return

    # Генерируем токен
    from auth_bom import generate_token
    token = generate_token(user.id)
    url = f"{BOM_SERVER_URL}/bom_login?token={token}"

    await update.message.reply_text(
        f"📋 <b>Состав продукции</b>\n\n"
        f"Ваша персональная ссылка (действует 7 дней):\n"
        f"<a href=\"{url}\">Открыть отчёт BOM</a>\n\n"
        f"<i>Ссылка привязана к вашему аккаунту.</i>",
        parse_mode='HTML',
        disable_web_page_preview=True
    )
    logger.info(f"BOM ссылка выдана пользователю {user.id} ({user.first_name})")


# ============================================================
# УПРАВЛЕНИЕ ПРАВИЛАМИ ФИЛЬТРАЦИИ (/rules)
# ============================================================
 
async def rules_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Показывает pending правила фильтрации для одобрения/отклонения."""
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return
 
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            # Pending правила
            cur.execute("""
                SELECT id, rule_type, value, reason, added_by, created_at
                FROM km_filter_rules
                WHERE approval_status = 'pending' AND is_active = true
                ORDER BY created_at DESC
                LIMIT 20
            """)
            pending = cur.fetchall()
 
            # Статистика
            cur.execute("""
                SELECT 
                    COUNT(*) FILTER (WHERE approval_status = 'approved' AND is_active = true) as approved,
                    COUNT(*) FILTER (WHERE approval_status = 'pending' AND is_active = true) as pending,
                    COUNT(*) FILTER (WHERE is_active = false) as disabled
                FROM km_filter_rules
            """)
            stats = cur.fetchone()
    finally:
        conn.close()
 
    if not pending:
        await update.message.reply_text(
            f"✅ Нет правил на рассмотрении\n\n"
            f"📊 Активных: {stats[0]}, отключённых: {stats[2]}\n\n"
            f"Поиск: /rules_find <слово>\n"
            f"Отключить: /rules_off <id>"
        )
        return
 
    await update.message.reply_text(
        f"📏 *Правила на рассмотрении: {len(pending)}*\n"
        f"Активных: {stats[0]} | Отключённых: {stats[2]}\n\n"
        f"Нажмите ✅ чтобы одобрить или ❌ чтобы отклонить:",
        parse_mode="Markdown"
    )
 
    for rule_id, rule_type, value, reason, added_by, created_at in pending:
        date_str = created_at.strftime("%d.%m") if created_at else ""
        text = (
            f"🏷 `{value}`\n"
            f"Тип: {rule_type} | От: {added_by or '?'} | {date_str}\n"
            f"Причина: {reason or '—'}"
        )
        keyboard = InlineKeyboardMarkup([
            [
                InlineKeyboardButton("✅ Одобрить", callback_data=f"rule_approve:{rule_id}"),
                InlineKeyboardButton("❌ Отклонить", callback_data=f"rule_reject:{rule_id}"),
            ]
        ])
        await update.message.reply_text(text, reply_markup=keyboard, parse_mode="Markdown")
 
 
async def rules_callback_handler(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик кнопок одобрения/отклонения правил."""
    query = update.callback_query
    await query.answer()
 
    if query.from_user.id != ADMIN_USER_ID:
        await query.answer("⛔ Только для администратора", show_alert=True)
        return
 
    data = query.data
 
    if data.startswith("rule_approve:"):
        rule_id = int(data.split(":")[1])
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    UPDATE km_filter_rules
                    SET approval_status = 'approved', updated_at = NOW()
                    WHERE id = %s
                    RETURNING value
                """, (rule_id,))
                row = cur.fetchone()
                conn.commit()
        finally:
            conn.close()
 
        value = row[0] if row else "?"
        await query.edit_message_text(f"✅ Правило одобрено: `{value}`", parse_mode="Markdown")
 
    elif data.startswith("rule_reject:"):
        rule_id = int(data.split(":")[1])
        conn = get_db_connection()
        try:
            with conn.cursor() as cur:
                cur.execute("""
                    UPDATE km_filter_rules
                    SET is_active = false, approval_status = 'rejected', updated_at = NOW()
                    WHERE id = %s
                    RETURNING value
                """, (rule_id,))
                row = cur.fetchone()
                conn.commit()
        finally:
            conn.close()
 
        value = row[0] if row else "?"
        await query.edit_message_text(f"❌ Правило отклонено: `{value}`", parse_mode="Markdown")


async def rules_find_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Поиск активных правил фильтрации по значению."""
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return

    if not context.args:
        await update.message.reply_text("Использование: /rules_find <слово>")
        return

    query_str = " ".join(context.args).strip()
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT id, value, rule_type, added_by, created_at::date
                FROM km_filter_rules
                WHERE is_active = true AND value ILIKE %s
                ORDER BY created_at DESC LIMIT 10
            """, (f"%{query_str}%",))
            rows = cur.fetchall()
    finally:
        conn.close()

    if not rows:
        await update.message.reply_text(f'🔍 Правила с "{query_str}" не найдены')
        return

    lines = [f'🔍 Найдено {len(rows)} для "{query_str}":\n']
    for i, (rid, value, rtype, added_by, created) in enumerate(rows, 1):
        date_str = created.strftime("%d.%m") if created else ""
        lines.append(f'{i}. [ID:{rid}] "{value}" ({rtype}) — {added_by or "?"}, {date_str}')
    all_ids = " ".join(str(r[0]) for r in rows)
    lines.append(f"\nОтключить: /rules_off {all_ids}")
    await update.message.reply_text("\n".join(lines))


async def rules_off_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Деактивация правил фильтрации по ID."""
    if update.effective_user.id != ADMIN_USER_ID:
        await update.message.reply_text("⛔ Только для администраторов")
        return

    if not context.args:
        await update.message.reply_text("Использование: /rules_off <id> [id2] [id3]...")
        return

    try:
        ids = [int(x) for x in context.args]
    except ValueError:
        await update.message.reply_text("ID должны быть числами")
        return

    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                UPDATE km_filter_rules SET is_active = false, updated_at = NOW()
                WHERE id = ANY(%s) AND is_active = true
                RETURNING id, value, rule_type
            """, (ids,))
            updated = cur.fetchall()
            conn.commit()
    finally:
        conn.close()

    if not updated:
        await update.message.reply_text("Ничего не отключено (правила не найдены или уже неактивны)")
        return

    lines = [f"✅ Отключено: {len(updated)}"]
    for rid, value, rtype in updated:
        lines.append(f'  "{value}" ({rtype}, ID:{rid})')
    await update.message.reply_text("\n".join(lines))

def get_chat_id_by_title(chat_title: str) -> int | None:
    """Получает chat_id по названию чата из БД."""
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute("""
                SELECT chat_id FROM tg_chats_metadata
                WHERE chat_title = %s
                ORDER BY last_message_at DESC
                LIMIT 1
            """, (chat_title,))
            result = cur.fetchone()
            if result:
                return result[0]
    except Exception as e:
        logger.error(f"Ошибка получения chat_id для '{chat_title}': {e}")
    finally:
        conn.close()
    return None


async def scheduled_daily_analysis(application):
    """Запланированный анализ документов в конце дня."""
    chat_id = get_chat_id_by_title(DELAYED_ANALYSIS_CHAT)

    if chat_id:
        await analyze_daily_documents(application.bot, chat_id, DELAYED_ANALYSIS_CHAT)
    else:
        logger.warning(f"Чат '{DELAYED_ANALYSIS_CHAT}' не найден для анализа документов")

async def handle_full_analysis_button(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик кнопки 'Полный анализ' — отправляет полный анализ в личку."""
    query = update.callback_query
    await query.answer()
    
    # Получаем message_id из callback_data
    callback_data = query.data
    if not callback_data.startswith("full_"):
        return
    
    try:
        original_message_id = int(callback_data.replace("full_", ""))
    except ValueError:
        await query.answer("Ошибка: неверный формат данных", show_alert=True)
        return
    
    chat_id = query.message.chat.id
    user_id = query.from_user.id
    chat_title = query.message.chat.title or "Чат"
    
    # Находим таблицу чата
    table_name = sanitize_table_name(chat_id, chat_title)
    
    # Получаем полный анализ из БД
    conn = get_db_connection()
    try:
        with conn.cursor() as cur:
            cur.execute(sql.SQL("""
                SELECT media_analysis, message_type, first_name
                FROM {}
                WHERE message_id = %s
            """).format(sql.Identifier(table_name)), (original_message_id,))
            
            result = cur.fetchone()
    finally:
        conn.close()
    
    if not result or not result[0]:
        await query.answer("Анализ не найден", show_alert=True)
        return
    
    media_analysis, message_type, sender_name = result
    
    # Формируем полное сообщение
    full_message = (
        f"📄 Полный анализ документа\n\n"
        f"📍 Чат: {chat_title}\n"
        f"👤 Отправил: {sender_name or 'Неизвестный'}\n"
        f"📎 Тип: {message_type}\n\n"
        f"{media_analysis.replace('*', '✱').replace('_', '‗')}"
    )
    
    # Отправляем в личку пользователю
    try:
        if len(full_message) > 4000:
            parts = [full_message[i:i+4000] for i in range(0, len(full_message), 4000)]
            for i, part in enumerate(parts):
                await context.bot.send_message(
                    chat_id=user_id,
                    text=part if i == 0 else f"...продолжение:\n\n{part}"
                )
        else:
            await context.bot.send_message(
                chat_id=user_id,
                text=full_message
            )
        await query.answer("✅ Полный анализ отправлен в личные сообщения")
        try:
            conn = get_db_connection()
            with conn.cursor() as cur:
                cur.execute(
                    "INSERT INTO bot_button_log (user_id, button_type, context_data) VALUES (%s, %s, %s)",
                    (user_id, 'full_analysis', json.dumps({"message_id": original_message_id, "chat_id": chat_id})),
                )
            conn.commit()
            conn.close()
        except Exception:
            pass
    except Exception as e:
        if "bot can't initiate" in str(e).lower() or "chat not found" in str(e).lower():
            await query.answer("❌ Сначала напишите боту в личку /start", show_alert=True)
        else:
            logger.error(f"Ошибка отправки полного анализа: {e}")
            await query.answer("❌ Ошибка отправки", show_alert=True)

async def element_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Обработчик команды /element — выдаёт данные для входа в Element X."""
    user_id = update.effective_user.id
    
    # Только в личных сообщениях
    if update.effective_chat.type != "private":
        await update.message.reply_text(
            "Напишите мне /element в личных сообщениях, и я отправлю вам данные для входа в Element X."
        )
        return
    
    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute(
            "SELECT matrix_username, matrix_password, joined_at FROM matrix_user_mapping WHERE telegram_user_id = %s",
            (user_id,)
        )
        row = cur.fetchone()
        cur.close()
        conn.close()
        
        if not row:
            await update.message.reply_text(
                "Ваш аккаунт Element X пока не создан. Обратитесь к администратору."
            )
            return
        
        mx_user, mx_pass, joined_at = row
        
        status = "✅ Вы уже подключились!" if joined_at else "⏳ Ожидает подключения"
        
        await update.message.reply_text(
            f"🔐 Ваши данные для входа в Element X\n"
            f"━━━━━━━━━━━━━━━\n"
            f"Сервер: frumelad.ru\n"
            f"Логин: {mx_user}\n"
            f"Пароль: {mx_pass}\n"
            f"Статус: {status}\n"
            f"━━━━━━━━━━━━━━━\n\n"
            f"📱 Как подключиться:\n"
            f"1. Скачайте Element X (не Element!):\n"
            f"   • iPhone: apps.apple.com/app/element-x/id1672254904\n"
            f"   • Android: play.google.com/store/apps/details?id=io.element.android.x\n"
            f"2. Откройте приложение\n"
            f"3. Нажмите «Изменить сервер» и введите: frumelad.ru\n"
            f"4. Введите логин и пароль выше\n"
            f"5. Все рабочие чаты доступны в пространстве «Фрумелад»\n\n"
            f"Сообщения синхронизируются с Telegram — можно пользоваться обоими мессенджерами."
        )
        
        # Отмечаем что инвайт отправлен
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute(
            "UPDATE matrix_user_mapping SET invited_at = NOW() WHERE telegram_user_id = %s AND invited_at IS NULL",
            (user_id,)
        )
        conn.commit()
        cur.close()
        conn.close()
        
    except Exception as e:
        logger.error(f"element_command error: {e}")
        await update.message.reply_text("Произошла ошибка. Попробуйте позже.")


async def rooms_command(update: Update, context: ContextTypes.DEFAULT_TYPE):
    """Команда /rooms — переотправить приглашения в Matrix-комнаты."""
    user_id = update.effective_user.id

    if update.effective_chat.type != "private":
        await update.message.reply_text("Напишите /rooms в личных сообщениях.")
        return

    import requests as req

    MATRIX_URL = os.environ.get("MATRIX_URL", "http://localhost:8008")
    MATRIX_ADMIN_USER = os.environ["MATRIX_ADMIN_USER"]
    MATRIX_ADMIN_PASSWORD = os.environ["MATRIX_ADMIN_PASSWORD"]
    SPACE_ROOM_ID = "!hRnxoPZwyiPRobHsCy:frumelad.ru"
    TG_TO_MATRIX_NAME = {
        "Руководство": "Руководство (bridged)",
        "Фрумелад (НБ) Кадровые задачи по IT и 1С": "Фрумелад (НФ) Кадровые задачи по IT и 1С",
    }
    WORK_ROOMS = {
        "Бухгалтерия Фрумелад/НФ", "Руководство (bridged)", "Производство",
        "Априори & Фрумелад/НФ", "Секретариат", "HR-Фрумелад/НФ",
        "Фрумелад задачи на разработку BSG", "Торты Отгрузки",
        "Фрумелад поддержка BSG", "Дизайн упаковки Кондитерская Прохорова",
        "Новые продукты и конкуренты", "БЗ Производство Chat", "БЗ Производство",
        "БЗ R&D", "БЗ R&D Chat", "БЗ Бухгалтерия", "БЗ Бухгалтерия Chat",
        "БЗ Закупки Chat", "БЗ Склад", "БЗ Склад Chat",
        "Подбор Персонала Внешний", "Отчеты по аутсорсингу",
        "R&D ~ общая рабочая группа",
        "KELIN - кондитерская Прохорова", "БЗ инструкции производство",
        "Закупки", "Закупки - Упаковка", "Продажи на ярды",
        "Склад - Производство",
        "Фрумелад (НФ) Кадровые задачи по IT и 1С",
        "Производство Кондитерская Прохорова", "HR Фрумелад",
    }

    try:
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute(
            "SELECT matrix_id FROM matrix_user_mapping WHERE telegram_user_id = %s",
            (user_id,)
        )
        row = cur.fetchone()
        if not row:
            await update.message.reply_text("Ваш аккаунт Element X не найден. Обратитесь к администратору.")
            cur.close(); conn.close()
            return
        matrix_id = row[0]

        # Get user's TG chats
        cur.execute("""
            SELECT m.chat_title FROM tg_user_roles ur
            JOIN tg_chats_metadata m ON m.chat_id = ur.chat_id
            WHERE ur.user_id = %s AND ur.is_active = true
        """, (user_id,))
        tg_chats = [r[0] for r in cur.fetchall()]
        cur.close(); conn.close()

        # Matrix login
        ms = req.Session()
        ms.trust_env = False
        login_resp = ms.post(f"{MATRIX_URL}/_matrix/client/v3/login", json={
            "type": "m.login.password",
            "user": MATRIX_ADMIN_USER,
            "password": MATRIX_ADMIN_PASSWORD,
        }, timeout=10).json()
        if "access_token" not in login_resp:
            await update.message.reply_text("Ошибка подключения к Matrix. Попробуйте позже.")
            return
        token = login_resp["access_token"]
        mh = {"Authorization": f"Bearer {token}"}

        # Get rooms
        room_map = {}
        _from = 0
        while True:
            rr = ms.get(f"{MATRIX_URL}/_synapse/admin/v1/rooms", headers=mh,
                        params={"limit": 100, "from": _from}, timeout=15).json()
            for r in rr.get("rooms", []):
                name = r.get("name")
                if name and name in WORK_ROOMS:
                    room_map[name] = r["room_id"]
            if len(rr.get("rooms", [])) < 100:
                break
            _from += 100

        # Space invite
        space_members = set(ms.get(
            f"{MATRIX_URL}/_synapse/admin/v1/rooms/{SPACE_ROOM_ID}/members",
            headers=mh, timeout=10
        ).json().get("members", []))

        invited_rooms = []
        already_in = []

        if matrix_id not in space_members:
            ms.post(f"{MATRIX_URL}/_matrix/client/v3/rooms/{SPACE_ROOM_ID}/invite",
                    headers=mh, json={"user_id": matrix_id}, timeout=10)
            invited_rooms.append("Пространство «Фрумелад»")

        # Room invites
        for chat_title in tg_chats:
            mn = TG_TO_MATRIX_NAME.get(chat_title, chat_title)
            if mn not in room_map:
                continue
            rid = room_map[mn]
            members = set(ms.get(
                f"{MATRIX_URL}/_synapse/admin/v1/rooms/{rid}/members",
                headers=mh, timeout=10
            ).json().get("members", []))
            if matrix_id in members:
                already_in.append(mn)
                continue
            ms.post(f"{MATRIX_URL}/_matrix/client/v3/rooms/{rid}/invite",
                    headers=mh, json={"user_id": matrix_id}, timeout=10)
            invited_rooms.append(mn)
            import time as _time; _time.sleep(0.2)

        if invited_rooms:
            rooms_list = "\n".join(f"  • {r}" for r in invited_rooms)
            msg = (
                f"✅ Приглашения отправлены!\n\n"
                f"Новые приглашения:\n{rooms_list}\n\n"
                f"Откройте Element X → раздел «Приглашения» и примите их."
            )
            if already_in:
                msg += f"\n\nВы уже состоите в: {', '.join(already_in[:5])}"
                if len(already_in) > 5:
                    msg += f" и ещё {len(already_in) - 5}"
        else:
            msg = "✅ Вы уже состоите во всех рабочих комнатах! Ничего отправлять не нужно."

        await update.message.reply_text(msg)

    except Exception as e:
        logger.error(f"rooms_command error: {e}", exc_info=True)
        await update.message.reply_text("Произошла ошибка. Попробуйте позже.")


async def element_reminder(context: ContextTypes.DEFAULT_TYPE):
    """Ежедневная проверка миграции на Element X.

    1. Синхронизирует joined_at через Synapse devices API.
    2. Проверяет membership в Space и комнатах для вошедших.
    3. Шлёт персональные напоминания в личку бота.
    4. Шлёт сводный отчёт админу.
    """
    import requests as req

    MATRIX_URL = os.environ.get("MATRIX_URL", "http://localhost:8008")
    MATRIX_ADMIN_USER = os.environ["MATRIX_ADMIN_USER"]
    MATRIX_ADMIN_PASSWORD = os.environ["MATRIX_ADMIN_PASSWORD"]
    SPACE_ROOM_ID = "!hRnxoPZwyiPRobHsCy:frumelad.ru"
    SKIP_MATRIX = {"@bot:frumelad.ru", "@aleksei:frumelad.ru"}
    TG_TO_MATRIX_NAME = {
        "Руководство": "Руководство (bridged)",
        "Фрумелад (НБ) Кадровые задачи по IT и 1С": "Фрумелад (НФ) Кадровые задачи по IT и 1С",
    }
    WORK_ROOMS = {
        "Бухгалтерия Фрумелад/НФ", "Руководство (bridged)", "Производство",
        "Априори & Фрумелад/НФ", "Секретариат", "HR-Фрумелад/НФ",
        "Фрумелад задачи на разработку BSG", "Торты Отгрузки",
        "Фрумелад поддержка BSG", "Дизайн упаковки Кондитерская Прохорова",
        "Новые продукты и конкуренты", "БЗ Производство Chat", "БЗ Производство",
        "БЗ R&D", "БЗ R&D Chat", "БЗ Бухгалтерия", "БЗ Бухгалтерия Chat",
        "БЗ Закупки Chat", "БЗ Склад", "БЗ Склад Chat",
        "Подбор Персонала Внешний", "Отчеты по аутсорсингу",
        "R&D ~ общая рабочая группа",
        "KELIN - кондитерская Прохорова", "БЗ инструкции производство",
        "Закупки", "Закупки - Упаковка", "Продажи на ярды",
        "Склад - Производство",
        "Фрумелад (НФ) Кадровые задачи по IT и 1С",
        "Производство Кондитерская Прохорова", "HR Фрумелад",
    }

    try:
        from proxy_config import get_proxy_url
        proxy_url = get_proxy_url()
        proxies = {"https": proxy_url, "http": proxy_url}

        # ── Matrix login ──
        ms = req.Session()
        ms.trust_env = False
        login_resp = ms.post(f"{MATRIX_URL}/_matrix/client/v3/login", json={
            "type": "m.login.password",
            "user": MATRIX_ADMIN_USER,
            "password": MATRIX_ADMIN_PASSWORD,
        }, timeout=10).json()
        if "access_token" not in login_resp:
            logger.error(f"element_reminder: Matrix login failed: {login_resp}")
            return
        token = login_resp["access_token"]
        mh = {"Authorization": f"Bearer {token}"}

        # ── Matrix rooms map ──
        room_map = {}
        _from = 0
        while True:
            rr = ms.get(f"{MATRIX_URL}/_synapse/admin/v1/rooms", headers=mh,
                        params={"limit": 100, "from": _from}, timeout=15).json()
            for r in rr.get("rooms", []):
                name = r.get("name")
                if name and name in WORK_ROOMS:
                    room_map[name] = r["room_id"]
            if len(rr.get("rooms", [])) < 100:
                break
            _from += 100

        space_members = set(ms.get(
            f"{MATRIX_URL}/_synapse/admin/v1/rooms/{SPACE_ROOM_ID}/members",
            headers=mh, timeout=10
        ).json().get("members", []))

        # ── Load mapping from DB ──
        conn = get_db_connection()
        cur = conn.cursor()
        cur.execute("""
            SELECT m.telegram_user_id, m.telegram_name, m.telegram_username,
                   m.matrix_id, m.matrix_username, m.matrix_password, m.joined_at
            FROM matrix_user_mapping m
        """)
        users = []
        for row in cur.fetchall():
            users.append({
                "tg_uid": row[0], "name": row[1], "username": row[2],
                "matrix_id": row[3], "mx_user": row[4], "mx_pass": row[5],
                "joined_at": row[6],
            })

        # ── Sync joined_at via devices API ──
        newly_joined = []
        for u in users:
            if u["matrix_id"] in SKIP_MATRIX:
                continue
            devs = ms.get(
                f"{MATRIX_URL}/_synapse/admin/v2/users/{u['matrix_id']}/devices",
                headers=mh, timeout=10
            ).json()
            has_devices = len(devs.get("devices", [])) > 0
            if has_devices and not u["joined_at"]:
                cur.execute(
                    "UPDATE matrix_user_mapping SET joined_at = NOW() WHERE telegram_user_id = %s",
                    (u["tg_uid"],)
                )
                u["joined_at"] = True  # mark locally
                newly_joined.append(u["name"])
            u["has_devices"] = has_devices
        conn.commit()

        # ── Get TG chat memberships ──
        for u in users:
            cur.execute("""
                SELECT m.chat_title FROM tg_user_roles ur
                JOIN tg_chats_metadata m ON m.chat_id = ur.chat_id
                WHERE ur.user_id = %s AND ur.is_active = true
            """, (u["tg_uid"],))
            u["tg_chats"] = [r[0] for r in cur.fetchall()]
        cur.close()
        conn.close()

        # ── Helpers: membership check & invite ──
        def get_membership(room_id, matrix_user_id):
            """Return membership state: 'join', 'invite', 'leave', 'ban', or None."""
            try:
                r = ms.get(
                    f"{MATRIX_URL}/_matrix/client/v3/rooms/{room_id}/state/m.room.member/{matrix_user_id}",
                    headers=mh, timeout=10)
                if r.status_code == 200:
                    return r.json().get("membership")
            except Exception:
                pass
            return None

        def send_matrix_invite(room_id, matrix_user_id):
            """Send Matrix room invite. Returns True on success."""
            r = ms.post(
                f"{MATRIX_URL}/_matrix/client/v3/rooms/{room_id}/invite",
                headers=mh, json={"user_id": matrix_user_id}, timeout=10)
            return r.status_code in (200, 403)

        # ── Check room membership for logged-in users ──
        room_members_cache = {}
        invites_sent = 0
        import time as _time
        for u in users:
            if u["matrix_id"] in SKIP_MATRIX:
                continue
            u["in_space"] = u["matrix_id"] in space_members
            u["missing_rooms"] = []
            if not u.get("has_devices"):
                continue

            # Space invite if needed
            if not u["in_space"]:
                ms_state = get_membership(SPACE_ROOM_ID, u["matrix_id"])
                if ms_state not in ("join", "invite"):
                    send_matrix_invite(SPACE_ROOM_ID, u["matrix_id"])
                    invites_sent += 1
                    _time.sleep(0.2)

            for chat_title in u["tg_chats"]:
                mn = TG_TO_MATRIX_NAME.get(chat_title, chat_title)
                if mn not in room_map:
                    continue
                rid = room_map[mn]
                if rid not in room_members_cache:
                    room_members_cache[rid] = set(ms.get(
                        f"{MATRIX_URL}/_synapse/admin/v1/rooms/{rid}/members",
                        headers=mh, timeout=10
                    ).json().get("members", []))
                if u["matrix_id"] in room_members_cache[rid]:
                    continue

                # Not a member — check state, send invite if needed
                ms_state = get_membership(rid, u["matrix_id"])
                if ms_state == "invite":
                    u["missing_rooms"].append(mn)
                elif ms_state in ("leave", None):
                    ok = send_matrix_invite(rid, u["matrix_id"])
                    if ok:
                        u["missing_rooms"].append(mn)
                        invites_sent += 1
                    _time.sleep(0.2)

        # ── Send personal reminders ──
        sent_personal = 0

        for u in users:
            if u["matrix_id"] in SKIP_MATRIX:
                continue
            tg_uid = u["tg_uid"]
            name = u["name"]

            if not u.get("has_devices"):
                # Не вошёл — напоминание с инструкцией
                msg = (
                    f"👋 {name}, напоминаю о переходе на Element X!\n\n"
                    f"Ваш аккаунт готов. Для входа:\n"
                    f"1. Скачайте Element X:\n"
                    f"   • Android: play.google.com/store/apps/details?id=io.element.android.x\n"
                    f"   • iPhone: apps.apple.com/app/element-x/id1672254904\n"
                    f"2. Нажмите «Изменить сервер» → введите: frumelad.ru\n"
                    f"3. Логин: {u['mx_user']}\n"
                    f"4. Пароль: напишите /element — я отправлю\n\n"
                    f"Все рабочие чаты уже ждут вас в пространстве «Фрумелад»."
                )
                try:
                    resp = req.post(
                        f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage",
                        json={"chat_id": tg_uid, "text": msg, "disable_web_page_preview": True},
                        proxies=proxies, timeout=15
                    ).json()
                    if resp.get("ok"):
                        sent_personal += 1
                    else:
                        u["_cant_send"] = True
                        logger.info(f"element_reminder: не удалось отправить {name} ({tg_uid}): {resp.get('description', '')}")
                except Exception as e:
                    u["_cant_send"] = True
                    logger.warning(f"element_reminder: ошибка отправки {name}: {e}")

            elif u.get("missing_rooms"):
                # Вошёл, приглашения отправлены/висят — напоминаем
                rooms_list = "\n".join(f"  • {r}" for r in u["missing_rooms"])
                space_note = ""
                if not u.get("in_space"):
                    space_note = "\n\n⚠️ Также примите приглашение в пространство «Фрумелад» — в нём собраны все рабочие комнаты."
                msg = (
                    f"👋 {name}, вы уже в Element X — отлично!\n\n"
                    f"Вам отправлены приглашения в комнаты:\n"
                    f"{rooms_list}\n\n"
                    f"Откройте Element X → раздел «Приглашения» и примите их, "
                    f"чтобы не пропускать рабочие обсуждения.{space_note}\n\n"
                    f"Если приглашения не видно — напишите /rooms, и я отправлю заново."
                )
                try:
                    resp = req.post(
                        f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage",
                        json={"chat_id": tg_uid, "text": msg, "disable_web_page_preview": True},
                        proxies=proxies, timeout=15
                    ).json()
                    if resp.get("ok"):
                        sent_personal += 1
                    else:
                        u["_cant_send"] = True
                        logger.info(f"element_reminder: не удалось отправить {name} ({tg_uid}): {resp.get('description', '')}")
                except Exception as e:
                    u["_cant_send"] = True
                    logger.warning(f"element_reminder: ошибка отправки {name}: {e}")

        # ── Group messages with inline button for unreachable users ──
        cant_reach = [u for u in users if u.get("_cant_send") and u["matrix_id"] not in SKIP_MATRIX]
        group_msgs_sent = 0
        if cant_reach:
            # Group unreachable users by their TG chats
            from collections import defaultdict
            chat_to_users = defaultdict(list)
            conn2 = get_db_connection()
            cur2 = conn2.cursor()
            for u in cant_reach:
                cur2.execute(
                    "SELECT chat_id FROM tg_user_roles WHERE user_id = %s AND is_active = true",
                    (u["tg_uid"],)
                )
                for (cid,) in cur2.fetchall():
                    chat_to_users[cid].append(u)
            cur2.close()
            conn2.close()

            # Send one message per group with inline button
            sent_chats = set()
            for chat_id, group_users in chat_to_users.items():
                if chat_id in sent_chats:
                    continue
                # Deduplicate users across chats
                names_in_group = []
                for u in group_users:
                    if u.get("username"):
                        names_in_group.append(f"@{u['username']}")
                    else:
                        names_in_group.append(u["name"])

                keyboard = {"inline_keyboard": [[{
                    "text": "📱 Получить данные для Element X",
                    "url": "https://t.me/AI_FRUM_NF_bot?start=element"
                }]]}
                group_msg = (
                    f"📢 Element X — корпоративный мессенджер\n\n"
                    f"Ещё не подключились: {', '.join(names_in_group)}\n\n"
                    f"Нажмите кнопку ниже — я отправлю вам данные для входа в личном сообщении."
                )
                try:
                    resp = req.post(
                        f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage",
                        json={
                            "chat_id": chat_id,
                            "text": group_msg,
                            "reply_markup": keyboard,
                            "disable_web_page_preview": True,
                        },
                        proxies=proxies, timeout=15
                    ).json()
                    if resp.get("ok"):
                        group_msgs_sent += 1
                        sent_chats.add(chat_id)
                except Exception as e:
                    logger.warning(f"element_reminder: group msg to {chat_id} failed: {e}")

        # ── Admin report ──
        total = len([u for u in users if u["matrix_id"] not in SKIP_MATRIX])
        joined = len([u for u in users if u.get("has_devices") and u["matrix_id"] not in SKIP_MATRIX])
        not_joined = total - joined
        with_missing = len([u for u in users if u.get("has_devices") and u.get("missing_rooms") and u["matrix_id"] not in SKIP_MATRIX])
        all_ok = len([u for u in users if u.get("has_devices") and not u.get("missing_rooms") and u["matrix_id"] not in SKIP_MATRIX])

        report_lines = [
            f"📊 Element X — статус миграции",
            f"",
            f"✅ Подключились полностью: {all_ok}",
            f"⚠️ Подключились, не все комнаты: {with_missing}",
            f"❌ Не вошли: {not_joined}",
            f"📨 Личных напоминаний отправлено: {sent_personal}",
            f"🔗 Matrix-приглашений отправлено: {invites_sent}",
            f"💬 Сообщений в группы (с кнопкой): {group_msgs_sent}",
        ]

        if newly_joined:
            report_lines.append(f"\n🆕 Новые подключения: {', '.join(newly_joined)}")

        if with_missing > 0:
            report_lines.append(f"\n⚠️ Не все комнаты приняты:")
            for u in users:
                if u.get("has_devices") and u.get("missing_rooms") and u["matrix_id"] not in SKIP_MATRIX:
                    rooms_str = ", ".join(u["missing_rooms"])
                    report_lines.append(f"  • {u['name']}: {rooms_str}")

        if not_joined > 0:
            report_lines.append(f"\n❌ Не вошли ({not_joined}):")
            names = []
            for u in users:
                if not u.get("has_devices") and u["matrix_id"] not in SKIP_MATRIX:
                    n = f"@{u['username']}" if u.get("username") else u["name"]
                    names.append(n)
            report_lines.append(f"  {', '.join(names)}")

        admin_report = "\n".join(report_lines)
        req.post(
            f"https://api.telegram.org/bot{BOT_TOKEN}/sendMessage",
            json={"chat_id": ADMIN_USER_ID, "text": admin_report, "disable_web_page_preview": True},
            proxies=proxies, timeout=30
        )

        logger.info(f"element_reminder: total={total}, joined={joined}, not_joined={not_joined}, missing_rooms={with_missing}, sent={sent_personal}")

    except Exception as e:
        logger.error(f"element_reminder error: {e}", exc_info=True)


def main():
    """Запуск бота."""
    if not BOT_TOKEN:
        logger.error("BOT_TOKEN не установлен в .env!")
        return

    if not DB_PASSWORD:
        logger.error("DB_PASSWORD не установлен в .env!")
        return

    from telegram.request import HTTPXRequest
    from proxy_config import get_proxy_url
    _proxy = get_proxy_url()
    request = HTTPXRequest(
        read_timeout=120,
        write_timeout=120,
        connect_timeout=30,
        proxy=_proxy
    )
    get_updates_request = HTTPXRequest(
        read_timeout=120,
        write_timeout=120,
        connect_timeout=30,
        proxy=_proxy
    )
    application = Application.builder().token(BOT_TOKEN).request(request).get_updates_request(get_updates_request).build()

    # Инициализация планировщика для отложенного анализа документов
    
    scheduler = AsyncIOScheduler()
    
    # Добавляем задачу на 23:55 каждый день
    
    scheduler.add_job(
        scheduled_daily_analysis,
        CronTrigger(hour=23, minute=55),
        args=[application],
        id='daily_document_analysis',
        name='Ежедневный анализ документов для группы "Торты Отгрузки"',
        replace_existing=True
    )
    
    # Запуск планировщика после инициализации event loop
    async def post_init(app):
        scheduler.start()
        logger.info(f"🕐 Планировщик запущен. Анализ документов для '{DELAYED_ANALYSIS_CHAT}' будет проводиться в 23:55")
    
    application.post_init = post_init
    
    # Команды
    application.add_handler(CommandHandler("start", start_command))
    application.add_handler(CommandHandler("roles", roles_command))
    application.add_handler(CommandHandler("stats", stats_command))
    application.add_handler(CommandHandler("search", search_command))
    application.add_handler(CommandHandler("chats", chats_command))
    application.add_handler(CommandHandler("analysis", analysis_command))
    application.add_handler(CommandHandler("analysis_list", analysis_list_command))
    # Email Logger команды
    application.add_handler(CommandHandler("threads", open_threads_command))
    application.add_handler(CommandHandler("email_stats", email_stats_command))
    application.add_handler(CommandHandler("sync_status", sync_status_command))
    application.add_handler(CommandHandler("search_email", search_email_command))
    application.add_handler(CommandHandler("add_employee", add_employee_command))
    application.add_handler(CommandHandler("assign_email", assign_email_command))
    application.add_handler(CommandHandler("list_employees", list_employees_command))
    application.add_handler(CommandHandler("bom", bom_command))
    # Команда /rules — управление правилами фильтрации
    application.add_handler(CommandHandler("rules", rules_command))
    application.add_handler(CommandHandler("rules_find", rules_find_command))
    application.add_handler(CommandHandler("rules_off", rules_off_command))
    application.add_handler(CommandHandler("element", element_command))
    application.add_handler(CommandHandler("rooms", rooms_command))

    # Notifications
    application.add_handler(get_notify_conversation_handler())
    application.add_handler(CallbackQueryHandler(handle_ack, pattern=r'^notify_ack_\d+$'))
    application.add_handler(CommandHandler("notify_status", notify_status))
    application.add_handler(CommandHandler("notify_remind", notify_remind))

    # Callback для правил
    application.add_handler(CallbackQueryHandler(
        rules_callback_handler,
        pattern=r'^rule_'
    ))

    application.add_handler(MessageHandler(
        filters.Regex(r'^/emailthread_\d+'),
        show_email_thread_command
    ))

    application.add_handler(CallbackQueryHandler(
        email_callback_handler,
        pattern=r'^email_'
    ))

    # === Nutrition Bot: запрос БЖУ у технологов ===
    from nutrition_bot import handle_callback as nutrition_callback
    from nutrition_bot import handle_text_reply as nutrition_text_reply
    from nutrition_bot import handle_photo as nutrition_photo
    
    async def nutrition_callback_handler(update, context):
        await update.callback_query.answer()
        nutrition_callback(update.callback_query.to_dict())
    
    async def nutrition_message_handler(update, context):
        msg = update.message
        if not msg or not msg.from_user:
            return
        # Проверяем есть ли активный nutrition запрос у этого пользователя
        try:
            from nutrition_bot import get_db
            conn = get_db()
            cur = conn.cursor()
            cur.execute("""
                SELECT id FROM nutrition_requests 
                WHERE assigned_to = %s AND status IN ('awaiting_input', 'awaiting_reject_reason')
                LIMIT 1
            """, (msg.from_user.id,))
            has_active = cur.fetchone() is not None
            cur.close()
            conn.close()
            if not has_active:
                return
        except:
            return
        
        if msg.photo:
            nutrition_photo(msg.to_dict())
        elif msg.text:
            nutrition_text_reply(msg.to_dict())
    
    async def nutrition_photo_handler(update, context):
        msg = update.message
        if not msg or not msg.from_user:
            return
        
        # Проверяем — фото или документ-изображение
        has_photo = bool(msg.photo)
        has_image_doc = (msg.document and msg.document.mime_type 
                        and msg.document.mime_type.startswith('image/'))
        
        if not has_photo and not has_image_doc:
            return
        
        try:
            from nutrition_bot import get_db
            conn = get_db()
            cur = conn.cursor()
            cur.execute("""
                SELECT id FROM nutrition_requests 
                WHERE assigned_to = %s AND status = 'awaiting_photo'
                LIMIT 1
            """, (msg.from_user.id,))
            has_active = cur.fetchone() is not None
            cur.close()
            conn.close()
            if not has_active:
                return
        except:
            return
        
        # Если документ-изображение — конвертируем в формат photo
        msg_dict = msg.to_dict()
        if has_image_doc and not has_photo:
            msg_dict['photo'] = [{'file_id': msg.document.file_id, 'file_unique_id': msg.document.file_unique_id, 'width': 0, 'height': 0}]
        
        nutrition_photo(msg_dict)
    
    application.add_handler(CallbackQueryHandler(
        nutrition_callback_handler,
        pattern=r'^nutr_'
    ))
    application.add_handler(MessageHandler(
        (filters.PHOTO | filters.Document.IMAGE) & filters.ChatType.PRIVATE,
        nutrition_photo_handler
    ), group=1)
    application.add_handler(MessageHandler(
        filters.TEXT & filters.ChatType.PRIVATE & ~filters.COMMAND & filters.REPLY,
        nutrition_message_handler
    ), group=1)
    
    # RAG агент в ЛИЧНЫХ сообщениях (без @упоминания)
    application.add_handler(MessageHandler(
        filters.TEXT & filters.ChatType.PRIVATE & ~filters.COMMAND & ~filters.REPLY,
        handle_private_rag
    ))

    # RAG агент — обработка упоминаний бота (ПЕРЕД log_message!)
    application.add_handler(MessageHandler(
        filters.TEXT & filters.Regex(r'@\w+'),
        handle_mention
    ))
    
    # Обработчик ответов с ролями (в группах и личных сообщениях)
    application.add_handler(MessageHandler(
        filters.TEXT & filters.REPLY & ~filters.COMMAND,
        handle_role_assignment
    ))
    
    # Обработчик всех остальных сообщений
    application.add_handler(MessageHandler(
        filters.ALL & ~filters.COMMAND,
        log_message
    ))

    # Обработчик кнопки "Полный анализ"
    application.add_handler(CallbackQueryHandler(handle_full_analysis_button))
    
    logger.info("🚀 Бот запущен. Логирование + анализ медиа + роли активны.")
    application.add_error_handler(error_handler)
    # Ежедневное напоминание о переходе на Element X (в 10:00 MSK = 07:00 UTC)
    from datetime import time as dt_time
    application.job_queue.run_daily(element_reminder, time=dt_time(hour=7, minute=0))
    application.run_polling(allowed_updates=Update.ALL_TYPES)


if __name__ == "__main__":
    main()
